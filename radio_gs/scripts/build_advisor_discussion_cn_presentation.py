#!/usr/bin/env python3
"""Build a Chinese advisor-discussion deck for the GaussFM paper."""

from __future__ import annotations

from pathlib import Path
from typing import Iterable, Sequence

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_AUTO_SIZE, MSO_VERTICAL_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


REPO_ROOT = Path(__file__).resolve().parents[2]
PAPER_DIR = REPO_ROOT / "paper"
FIG_DIR = PAPER_DIR / "figures"
OUT_PPTX = PAPER_DIR / "gaussfm_advisor_discussion_cn.pptx"
OUT_OUTLINE = PAPER_DIR / "gaussfm_advisor_discussion_cn_outline.md"
FRAMEWORK_FIG = FIG_DIR / "gaussfm_overall_framework_latest.png"
MEMORY_FIG = FIG_DIR / "gaussfm_compact_gaussian_feature_memory_latest.png"
QUAL_LERF = FIG_DIR / "lerf_2d3d_ovs_qualitative.png"
QUAL_SCAN = FIG_DIR / "scannet_openvocab_3d_query_qualitative_gaussfm.png"
QUAL_RADIO = FIG_DIR / "lerf_rendered_grounding_radio_vs_gaussfm.png"
QUAL_SAM_DINO = FIG_DIR / "lerf_sam_dino_tasks_qualitative_gaussfm.png"
QUAL_DIRECT = FIG_DIR / "lerf_main_qualitative_comparison_gaussfm.png"

SLIDE_W = 13.333
SLIDE_H = 7.5
FONT = "Microsoft YaHei"

COLORS = {
    "ink": RGBColor(23, 31, 42),
    "muted": RGBColor(82, 94, 109),
    "paper": RGBColor(248, 250, 252),
    "panel": RGBColor(255, 255, 255),
    "line": RGBColor(214, 220, 228),
    "teal": RGBColor(18, 126, 126),
    "blue": RGBColor(44, 93, 165),
    "green": RGBColor(68, 135, 93),
    "amber": RGBColor(196, 123, 39),
    "red": RGBColor(174, 68, 68),
    "violet": RGBColor(95, 88, 170),
    "soft_teal": RGBColor(229, 246, 245),
    "soft_blue": RGBColor(232, 239, 252),
    "soft_amber": RGBColor(252, 244, 230),
    "soft_red": RGBColor(252, 236, 236),
}


class AdvisorDeck:
    def __init__(self) -> None:
        self.prs = Presentation()
        self.prs.slide_width = Inches(SLIDE_W)
        self.prs.slide_height = Inches(SLIDE_H)
        self.blank = self.prs.slide_layouts[6]
        self.outline: list[tuple[str, list[str]]] = []

    def save(self) -> None:
        OUT_PPTX.parent.mkdir(parents=True, exist_ok=True)
        self.prs.save(OUT_PPTX)
        lines = [
            "# GaussFM 中文导师讨论版 PPT 大纲",
            "",
            f"- PPTX: `{OUT_PPTX.relative_to(REPO_ROOT)}`",
            "- 口径：标准会议论文 presentation 流程 / 中文导师讨论版 / 投稿故事线",
            "- 结构：背景 -> 动机 -> 相关工作 -> 问题缺口 -> 方法 -> 实验设置 -> 主结果 -> 定性 -> 控制对比 -> 消融 -> 边界",
            "",
        ]
        for idx, (title, notes) in enumerate(self.outline, 1):
            lines.append(f"## {idx}. {title}")
            for note in notes:
                lines.append(f"- {note}")
            lines.append("")
        OUT_OUTLINE.write_text("\n".join(lines), encoding="utf-8")

    def slide(self, title: str, section: str | None = None) -> object:
        slide = self.prs.slides.add_slide(self.blank)
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = COLORS["paper"]
        if section:
            self.text(slide, 0.50, 0.18, 4.6, 0.24, section.upper(), 8, COLORS["teal"], bold=True)
        self.text(slide, 0.50, 0.42, 12.1, 0.50, title, 22, COLORS["ink"], bold=True)
        self.line(slide, 0.50, 1.03, 12.80, 1.03, COLORS["line"])
        self.footer(slide)
        self.outline.append((title, []))
        return slide

    def note(self, text: str) -> None:
        self.outline[-1][1].append(text)

    def footer(self, slide: object) -> None:
        idx = len(self.prs.slides)
        self.text(slide, 0.50, 7.08, 4.3, 0.18, "GaussFM 导师讨论版", 7, COLORS["muted"])
        self.text(slide, 12.15, 7.08, 0.55, 0.18, str(idx), 7, COLORS["muted"], align=PP_ALIGN.RIGHT)

    def text(
        self,
        slide: object,
        x: float,
        y: float,
        w: float,
        h: float,
        value: str,
        size: int,
        color: RGBColor | None = None,
        bold: bool = False,
        align: PP_ALIGN | None = None,
    ) -> object:
        box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = box.text_frame
        tf.clear()
        tf.word_wrap = True
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        tf.margin_left = Inches(0.02)
        tf.margin_right = Inches(0.02)
        tf.margin_top = Inches(0.01)
        tf.margin_bottom = Inches(0.01)
        p = tf.paragraphs[0]
        if align is not None:
            p.alignment = align
        r = p.add_run()
        r.text = value
        r.font.name = FONT
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color or COLORS["ink"]
        return box

    def bullets(
        self,
        slide: object,
        x: float,
        y: float,
        w: float,
        h: float,
        items: Sequence[str],
        size: int = 14,
        color: RGBColor | None = None,
    ) -> None:
        box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = box.text_frame
        tf.clear()
        tf.word_wrap = True
        tf.margin_left = Inches(0.10)
        tf.margin_right = Inches(0.04)
        tf.margin_top = Inches(0.02)
        tf.margin_bottom = Inches(0.02)
        for i, item in enumerate(items):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = item
            p.level = 0
            p.space_after = Pt(6)
            p.font.name = FONT
            p.font.size = Pt(size)
            p.font.color.rgb = color or COLORS["ink"]

    def card(
        self,
        slide: object,
        x: float,
        y: float,
        w: float,
        h: float,
        title: str,
        body: str,
        accent: RGBColor,
        body_size: int = 14,
    ) -> None:
        shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
        shape.fill.solid()
        shape.fill.fore_color.rgb = COLORS["panel"]
        shape.line.color.rgb = COLORS["line"]
        shape.adjustments[0] = 0.04
        bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x), Inches(y), Inches(0.06), Inches(h))
        bar.fill.solid()
        bar.fill.fore_color.rgb = accent
        bar.line.fill.background()
        self.text(slide, x + 0.17, y + 0.13, w - 0.28, 0.26, title, 10, COLORS["muted"], bold=True)
        self.text(slide, x + 0.17, y + 0.47, w - 0.28, h - 0.52, body, body_size, COLORS["ink"], bold=True)

    def chip(self, slide: object, x: float, y: float, w: float, label: str, color: RGBColor) -> None:
        shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(0.33))
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()
        shape.adjustments[0] = 0.5
        tf = shape.text_frame
        tf.clear()
        tf.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = label
        r.font.name = FONT
        r.font.size = Pt(8)
        r.font.bold = True
        r.font.color.rgb = RGBColor(255, 255, 255)

    def line(self, slide: object, x1: float, y1: float, x2: float, y2: float, color: RGBColor) -> None:
        line = slide.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
        line.line.color.rgb = color
        line.line.width = Pt(1.0)

    def image(self, slide: object, path: Path, x: float, y: float, w: float, h: float) -> None:
        if not path.exists():
            self.card(slide, x, y, w, h, "Missing figure", str(path.relative_to(REPO_ROOT)), COLORS["red"], 10)
            return
        with Image.open(path) as im:
            iw, ih = im.size
        box_ratio = w / h
        img_ratio = iw / ih
        if img_ratio >= box_ratio:
            draw_w = w
            draw_h = w / img_ratio
            draw_x = x
            draw_y = y + (h - draw_h) / 2
        else:
            draw_h = h
            draw_w = h * img_ratio
            draw_x = x + (w - draw_w) / 2
            draw_y = y
        slide.shapes.add_picture(str(path), Inches(draw_x), Inches(draw_y), Inches(draw_w), Inches(draw_h))

    def table(
        self,
        slide: object,
        x: float,
        y: float,
        w: float,
        h: float,
        data: Sequence[Sequence[str]],
        font_size: int = 9,
        emphasize_rows: Iterable[int] = (),
        emphasize_cols: Iterable[int] = (),
    ) -> None:
        rows, cols = len(data), len(data[0])
        shape = slide.shapes.add_table(rows, cols, Inches(x), Inches(y), Inches(w), Inches(h))
        table = shape.table
        emph_rows = set(emphasize_rows)
        emph_cols = set(emphasize_cols)
        for i, row in enumerate(data):
            for j, value in enumerate(row):
                cell = table.cell(i, j)
                cell.text = value
                cell.margin_left = Inches(0.03)
                cell.margin_right = Inches(0.03)
                cell.margin_top = Inches(0.02)
                cell.margin_bottom = Inches(0.02)
                if i == 0:
                    fill = COLORS["soft_blue"]
                elif i in emph_rows:
                    fill = COLORS["soft_teal"]
                elif j in emph_cols:
                    fill = COLORS["soft_amber"]
                else:
                    fill = COLORS["panel"]
                cell.fill.solid()
                cell.fill.fore_color.rgb = fill
                for p in cell.text_frame.paragraphs:
                    p.alignment = PP_ALIGN.CENTER if j else PP_ALIGN.LEFT
                    for r in p.runs:
                        r.font.name = FONT
                        r.font.size = Pt(font_size)
                        r.font.bold = i == 0 or i in emph_rows
                        r.font.color.rgb = COLORS["ink"]


def _build_deck_legacy() -> AdvisorDeck:
    d = AdvisorDeck()

    # 1
    s = d.slide("GaussFM：面向开放词汇三维理解的紧凑基础特征高斯记忆", "定位")
    d.text(s, 0.85, 1.40, 11.8, 0.80, "核心主张：把 frame-wise RADIO 特征转化为可部署、可渲染、可直接查询的 3D scene memory。", 24, COLORS["ink"], bold=True, align=PP_ALIGN.CENTER)
    d.card(s, 1.05, 2.75, 3.45, 1.25, "不是", "新的 RGB 3DGS 重建器", COLORS["red"], 15)
    d.card(s, 4.95, 2.75, 3.45, 1.25, "而是", "compact foundation-feature memory", COLORS["teal"], 15)
    d.card(s, 8.85, 2.75, 3.45, 1.25, "支撑", "2D / 3D / point-level open-vocabulary query", COLORS["blue"], 15)
    d.text(s, 1.20, 4.80, 11.0, 0.45, "导师讨论目标：确认故事线、实验充分性、以及投稿前还需要补哪些证据。", 17, COLORS["muted"], align=PP_ALIGN.CENTER)
    d.note("开场直接说明讨论目标，不做对外包装。")

    # 2
    s = d.slide("为什么 2D foundation features 不能直接成为 3D scene memory", "问题")
    d.card(s, 0.85, 1.35, 3.70, 1.45, "View-local", "单帧特征受视角、遮挡、尺度影响；跨视角不天然一致。", COLORS["blue"], 14)
    d.card(s, 4.82, 1.35, 3.70, 1.45, "High-dimensional", "直接存 1280-D RADIO 特征，场景规模一大就很重。", COLORS["amber"], 14)
    d.card(s, 8.80, 1.35, 3.70, 1.45, "Not directly 3D", "2D heatmap 不能直接回答 Gaussian primitive / point query。", COLORS["red"], 14)
    d.bullets(s, 1.10, 3.45, 11.1, 1.7, [
        "现有 open-vocabulary 3D 表示常在 rendered-view query、direct 3D query、point query 之间割裂。",
        "本文要解决的是：同一个紧凑三维特征记忆，能否支撑多种开放词汇查询协议。",
        "因此主线不是“压缩特征”本身，而是“紧凑存储 + 可重建 + 可查询 + 下游可用”。",
    ], 15)
    d.note("这页建立动机：3D 场景需要自己的 foundation-feature memory。")

    # 3
    s = d.slide("核心科学问题：高维、多视角、帧级特征如何变成三维记忆", "问题")
    d.text(s, 0.95, 1.35, 11.5, 0.75, "目标不是保存所有帧特征，而是在 3D Gaussian 上学习一个 compact RADIO-compatible feature field。", 21, COLORS["ink"], bold=True, align=PP_ALIGN.CENTER)
    d.card(s, 1.00, 2.75, 3.45, 1.35, "Q1 存储", "如何低维存储但保持 RADIO-space 可重建性？", COLORS["teal"], 14)
    d.card(s, 4.95, 2.75, 3.45, 1.35, "Q2 一致性", "如何把多视角证据压回稳定 primitive support？", COLORS["blue"], 14)
    d.card(s, 8.90, 2.75, 3.45, 1.35, "Q3 查询", "如何让同一记忆支持 2D、3D 和 point query？", COLORS["green"], 14)
    d.text(s, 1.20, 5.00, 10.9, 0.42, "论文成败取决于：方法组件是否对应这些问题，实验是否分别闭环。", 17, COLORS["red"], bold=True, align=PP_ALIGN.CENTER)
    d.note("该页也解释为什么后续实验分成三条主线。")

    # 4
    s = d.slide("方法总览：训练时重建 foundation features，推理时查询 compact memory", "方法")
    d.image(s, FIG_DIR / "figure1_overall_framework.png", 0.75, 1.20, 12.0, 4.85)
    d.text(s, 0.95, 6.15, 11.5, 0.42, "一句话：每个 Gaussian 存 compact latent；按需解码为 RADIO-compatible features；同一场景记忆服务三类 query。", 15, COLORS["muted"], align=PP_ALIGN.CENTER)
    d.note("导师讨论时可重点看：图是否过复杂、命名是否足够清晰。")

    # 5
    s = d.slide("核心组件一：Compact Gaussian Feature Field", "方法")
    d.image(s, FIG_DIR / "figure2_method_details.png", 0.70, 1.25, 6.20, 4.80)
    d.card(s, 7.10, 1.30, 5.10, 0.88, "低维 Gaussian code", "不存完整 1280-D RADIO feature，而存 compact latent。", COLORS["teal"], 13)
    d.card(s, 7.10, 2.35, 5.10, 0.88, "空间上下文", "通过 spatial context / neighborhood support 补足单个 Gaussian 的局部性。", COLORS["blue"], 13)
    d.card(s, 7.10, 3.40, 5.10, 0.88, "可靠性与可见性", "区分可见、高置信、多视角支持稳定的 primitive。", COLORS["green"], 13)
    d.text(s, 7.15, 4.85, 5.00, 0.70, "需要强调：compact 不是单纯降维，而是把多视角信息组织成场景级记忆。", 15, COLORS["ink"], bold=True)
    d.note("方法组件页 1：讲清 compact field 不是一个简单 MLP decoder。")

    # 6
    s = d.slide("核心组件二：Foundation-space reconstruction 与多视角支持", "方法")
    d.card(s, 0.90, 1.35, 3.75, 1.45, "HCD / CTR codec", "把 compact latent 解码回 RADIO-compatible feature space。", COLORS["teal"], 14)
    d.card(s, 4.80, 1.35, 3.75, 1.45, "screen-space refiner", "提高渲染 feature map 的像素对齐与边界质量。", COLORS["blue"], 14)
    d.card(s, 8.70, 1.35, 3.75, 1.45, "multiview registration", "把跨视角 primitive evidence 压回 compact field。", COLORS["green"], 14)
    d.bullets(s, 1.10, 3.30, 11.0, 1.70, [
        "训练：RADIO-space feature reconstruction + 多视角/几何/可见性约束。",
        "推理：不重新跑 frame-wise feature extraction；直接从存储的 Gaussian memory 读出。",
        "关键边界：VPR/冻结头用于训练或诊断支持，不应被写成部署时必须依赖的主路径。",
    ], 15)
    d.note("这页把训练机制和推理依赖边界说清楚。")

    # 7
    s = d.slide("主要贡献 / 创新点：三条 claim 必须和证据一一对应", "贡献")
    d.card(s, 0.80, 1.35, 3.85, 2.10, "贡献 1", "提出 compact foundation-feature Gaussian memory：低维存储、按需重建 RADIO-compatible features。", COLORS["teal"], 14)
    d.card(s, 4.75, 1.35, 3.85, 2.10, "贡献 2", "统一 2D rendered query、3D primitive query、ScanNet point query 三种开放词汇接口。", COLORS["blue"], 14)
    d.card(s, 8.70, 1.35, 3.85, 2.10, "贡献 3", "系统验证 compact field 相比 frame-wise RADIO / cache / 1280-D memory 的下游可用性和存储优势。", COLORS["green"], 14)
    d.text(s, 1.05, 4.35, 11.3, 0.75, "审稿风险点：如果只讲主结果、不讲 RADIO 对比和消融，方法会显得像工程拼接；所以后面实验要放全。", 18, COLORS["red"], bold=True, align=PP_ALIGN.CENTER)
    d.note("该页明确创新点和实验必要性。")

    # 8
    s = d.slide("实验矩阵：每个 claim 都要有直接证据", "实验总览")
    d.table(s, 0.75, 1.25, 11.85, 4.70, [
        ["Claim", "核心证据", "补充/边界"],
        ["2D rendered query 可用", "LERF rendered-view OVS", "feature-only boundary query interface"],
        ["3D primitive query 可用", "LERF direct 3D object selection", "pure compact / guard / SAM3 diagnostic 分开"],
        ["跨数据集 point query", "VALA-aligned ScanNet table", "contextual kNN / spatial smoothing diagnostic"],
        ["不是低质压缩", "frame-wise RADIO vs GaussFM", "nearest cache / 1280-D memory controls"],
        ["组件确实必要", "architecture + query interface ablations", "非所有模块单调提升，需诚实呈现"],
    ], 11, emphasize_rows=[1, 2, 3])
    d.note("该页作为后续实验页导航。")

    # 9
    s = d.slide("主结果一：LERF rendered-view open-vocabulary query", "主结果")
    d.table(s, 0.85, 1.30, 5.45, 2.70, [
        ["Method", "mIoU", "Acc"],
        ["LangSplat", "33.92", "64.01"],
        ["LangSplatV2", "46.24", "75.84"],
        ["OpenGaussian", "49.74", "72.75"],
        ["GaussFM", "64.98", "82.68"],
    ], 11, emphasize_rows=[4])
    d.card(s, 6.80, 1.30, 5.55, 1.15, "这页证明什么", "重建出的 dense feature map 能支撑 novel-view 2D open-vocabulary grounding。", COLORS["teal"], 15)
    d.card(s, 6.80, 2.78, 5.55, 1.15, "为什么重要", "说明 compact field 不只是存储压缩，而保留了文本查询可用性。", COLORS["blue"], 15)
    d.card(s, 6.80, 4.26, 5.55, 1.15, "讨论点", "是否需要把 feature-only boundary query interface 和核心 heatmap 能力更清楚地区分？", COLORS["amber"], 15)
    d.note("主结果 1：注意解释 mIoU/Acc 和 query interface 边界。")

    # 10
    s = d.slide("主结果二：LERF direct 3D object selection", "主结果")
    d.table(s, 0.85, 1.30, 5.60, 2.70, [
        ["Method", "mIoU", "Acc"],
        ["OpenGaussian", "41.06", "51.44"],
        ["Dr. Splat", "39.77", "65.48"],
        ["LangSplatV2", "35.87", "55.80"],
        ["GaussFM", "54.36", "80.84"],
    ], 11, emphasize_rows=[4])
    d.bullets(s, 6.90, 1.35, 5.45, 2.10, [
        "证明同一 Gaussian memory 可在 primitive/object 层面直接查询。",
        "与 rendered-view heatmap 是不同能力：不依赖先渲染 2D mask 再反投影。",
        "Direct3D 页必须明确 pure compact、score guard、SAM3 diagnostic 的边界。",
    ], 15)
    d.card(s, 6.90, 4.35, 5.35, 1.05, "导师讨论", "当前 Direct3D 最容易被质疑的是 support policy 是否像后处理；消融页要专门解释。", COLORS["red"], 14)
    d.note("主结果 2：强调 primitive-level query 和支持策略边界。")

    # 11
    s = d.slide("主结果三：VALA-aligned ScanNet point query", "主结果")
    d.table(s, 0.65, 1.25, 7.35, 3.45, [
        ["Method", "19 mIoU", "19 mAcc", "15 mIoU", "15 mAcc", "10 mIoU", "10 mAcc"],
        ["LangSplatV2", "14.75", "25.47", "17.09", "35.68", "22.83", "41.52"],
        ["OpenGaussian", "27.73", "42.01", "29.67", "46.15", "39.93", "57.34"],
        ["Dr. Splat", "29.31", "47.68", "33.25", "54.33", "44.19", "65.19"],
        ["VALA", "32.11", "50.05", "35.10", "54.77", "46.21", "65.61"],
        ["GaussFM", "36.55", "50.57", "42.78", "72.85", "57.85", "77.93"],
    ], 8, emphasize_rows=[5])
    d.card(s, 8.40, 1.45, 4.20, 1.30, "这页证明什么", "方法不是 LERF 小场景特化；可迁移到 ScanNet point-level open-vocabulary query。", COLORS["teal"], 13)
    d.card(s, 8.40, 3.10, 4.20, 1.30, "必须谨慎", "这是 VALA-aligned 8-scene direct point-query protocol，不写 full ScanNet semantic segmentation SOTA。", COLORS["amber"], 13)
    d.note("主结果 3：避免把 point-query probe 包装成完整 ScanNet leaderboard。")

    # 12
    s = d.slide("原始 RADIO 对比一：不是 nearest-view cache，也不是 1280-D 显式存储", "RADIO 对比")
    d.table(s, 0.62, 1.22, 8.20, 4.10, [
        ["Method", "Compact", "3D memory", "LERF LocAcc", "LERF mIoU", "Storage"],
        ["Frame-wise RADIO", "no", "no", "0.7985", "0.4634", "per-frame cache"],
        ["Nearest-view RADIO cache", "no", "no", "0.2722", "0.1545", "per-frame cache"],
        ["Per-Gaussian 1280-D RADIO", "no", "yes", "0.5642", "0.3182", "1039.7 MiB mean"],
        ["GaussFM core", "yes", "yes", "0.8598", "0.5707", "compact checkpoint"],
        ["GaussFM + boundary", "yes", "yes", "0.8598", "0.5889", "same memory + query interface"],
    ], 9, emphasize_rows=[4, 5])
    d.card(s, 9.20, 1.55, 3.25, 1.30, "结论", "GaussFM 优势不是来自最近帧 cache，也不是粗暴存 1280-D feature。", COLORS["teal"], 14)
    d.card(s, 9.20, 3.25, 3.25, 1.30, "讨论点", "这页应放进主文还是补充？若主文空间紧，至少 presentation 必须讲清。", COLORS["amber"], 14)
    d.note("原始 RADIO controlled evidence 是导师最关心的支撑页。")

    # 13
    s = d.slide("原始 RADIO 对比二：SAM3 / DINO frozen-head downstream tasks", "RADIO 对比")
    d.table(s, 0.50, 1.18, 12.30, 4.45, [
        ["Task", "Primary", "Frame-wise RADIO", "GaussFM", "Delta", "Secondary caveat"],
        ["SAM3 point prompt", "mIoU", "0.3700", "0.4173", "+0.0473", "LocAcc same"],
        ["SAM3 box prompt", "mIoU", "0.6560", "0.6638", "+0.0079", "LocAcc lower"],
        ["SAM3 mask propagation", "mIoU", "0.3583", "0.3756", "+0.0173", "LocAcc lower"],
        ["DINOv3 dense matching", "Mean score", "0.8547", "0.9048", "+0.0501", "HitRate lower"],
        ["DINOv3 mask propagation", "mIoU", "0.4606", "0.4677", "+0.0071", "LocAcc higher"],
    ], 9, emphasize_rows=[1, 2, 3, 4, 5])
    d.text(s, 0.75, 5.95, 11.85, 0.45, "Claim-safe wording：selected primary downstream feature-usability metrics 全部提升；不宣称 universal superiority。", 16, COLORS["red"], bold=True, align=PP_ALIGN.CENTER)
    d.note("该页把 6/6 primary wins 和 caveats 放在同一页，避免过度包装。")

    # 14
    s = d.slide("完整消融一：architecture ablation 显示 HCD/重建 codec 是最大依赖", "消融")
    d.table(s, 0.70, 1.20, 7.20, 3.95, [
        ["Variant", "LocAcc", "mIoU", "Reading"],
        ["Full seed-7 GaussFM", "0.8580", "0.4850", "reference"],
        ["w/o FDH warm-start", "0.8018", "0.4236", "几何/特征初始化重要"],
        ["w/o refiner", "0.8401", "0.4796", "主要影响较小"],
        ["w/o hybrid", "0.8394", "0.5069", "peak 与 region 有 tradeoff"],
        ["w/o HCD / CTR", "0.5306", "0.2596", "最大退化"],
    ], 10, emphasize_rows=[1])
    d.card(s, 8.35, 1.40, 4.00, 1.18, "关键结论", "Foundation-space reconstruction codec 不是小模块，是方法成立的核心。", COLORS["teal"], 14)
    d.card(s, 8.35, 2.90, 4.00, 1.18, "诚实边界", "Hybrid 并非所有指标单调提升，应解释为 peak stability / memory design tradeoff。", COLORS["amber"], 14)
    d.card(s, 8.35, 4.40, 4.00, 1.18, "导师讨论", "主文消融是否需要从“组件列表”改成“问题-组件-证据”组织？", COLORS["blue"], 14)
    d.note("architecture 消融比短讲版更完整。")

    # 15
    s = d.slide("完整消融二：Direct3D query interface / support calibration 需要分层汇报", "消融")
    d.table(s, 0.65, 1.20, 7.75, 4.10, [
        ["Query interface", "mIoU", "Acc@0.25", "Boundary-F", "Use"],
        ["strict single-prompt pure one-map", "0.4489", "0.6724", "0.6124", "strict compact"],
        ["pure compact prompt ensemble", "0.4570", "0.6851", "0.6166", "compact core"],
        ["score-component guard", "0.5014", "0.7044", "0.6305", "main compact"],
        ["official SAM3 box", "0.5705", "0.6835", "0.6681", "diagnostic"],
    ], 9, emphasize_rows=[3])
    d.card(s, 8.75, 1.50, 3.70, 1.20, "必须区分", "score-component guard 是主 compact query interface；official SAM3 box 是边界潜力 diagnostic。", COLORS["red"], 13)
    d.card(s, 8.75, 3.05, 3.70, 1.20, "审稿风险", "如果不分层，reviewer 会认为主要结果依赖外部后处理。", COLORS["amber"], 13)
    d.card(s, 8.75, 4.60, 3.70, 1.20, "补充建议", "可在 appendix 放 full query interface registry，主文只保留 compact row + diagnostic caveat。", COLORS["blue"], 13)
    d.note("Direct3D ablation 直接回应导师对 query interface 的疑虑。")

    # 16
    s = d.slide("完整消融三：ScanNet query interface diagnostic 说明 point-query 仍有调参/读出因素", "消融")
    d.table(s, 0.55, 1.20, 12.10, 4.55, [
        ["Variant", "19 mIoU/mAcc", "15 mIoU/mAcc", "10 mIoU/mAcc", "Use"],
        ["DINO-CV kNN alpha=0.5", "0.3704/0.6017", "0.3771/0.6198", "0.4585/0.7032", "support"],
        ["k12/cand48 alpha=0.45", "0.3715/0.6024", "0.3784/0.6206", "0.4585/0.7029", "balanced support"],
        ["k16/cand80 + spatial k12/a1", "0.3806/0.6129", "0.3871/0.6315", "0.4711/0.7200", "supporting diagnostic"],
        ["proposal memory", "0.3931/0.6255", "0.3837/0.6228", "0.4612/0.7081", "19-class detail only"],
        ["paper-facing row", "0.3655/0.5057", "0.4278/0.7285", "0.5785/0.7793", "main table"],
    ], 8, emphasize_rows=[5])
    d.text(s, 0.85, 5.95, 11.60, 0.42, "直接说清：ScanNet 里主表、support diagnostic、proposal-memory ablation 不是同一种 claim。", 15, COLORS["red"], bold=True, align=PP_ALIGN.CENTER)
    d.note("这页避免 ScanNet 证据源混淆，也是之前 validator 暴露的问题。")

    # 17
    s = d.slide("Storage / efficiency：应同时报告 latent、package、full checkpoint", "效率")
    d.table(s, 0.55, 1.15, 12.20, 4.15, [
        ["Scene", "Direct 1280-D", "Latent payload", "Latent saving", "Feature package", "Package saving", "Full saving"],
        ["Figurines", "412.1 MiB", "20.6 MiB", "20.00x", "199.0 MiB", "2.07x", "1.74x"],
        ["Ramen", "934.3 MiB", "46.7 MiB", "20.00x", "225.1 MiB", "4.15x", "3.00x"],
        ["Teatime", "1123.4 MiB", "56.2 MiB", "20.00x", "234.5 MiB", "4.79x", "3.32x"],
        ["Waldo Kitchen", "1688.8 MiB", "84.4 MiB", "20.00x", "262.8 MiB", "6.43x", "4.04x"],
    ], 8, emphasize_cols=[3])
    d.card(s, 0.95, 5.70, 5.65, 0.82, "汇报口径", "latent payload 是最干净的 feature-memory 压缩；package/full checkpoint 是部署视角。", COLORS["teal"], 13)
    d.card(s, 6.95, 5.70, 5.25, 0.82, "风险", "只报 full checkpoint 会低估 compact feature-memory 的真实优势。", COLORS["amber"], 13)
    d.note("这页用于和导师讨论论文表格该报哪几列。")

    # 18
    s = d.slide("定性证据：只服务核心 claim，不做图像堆砌", "定性")
    d.card(s, 0.85, 1.35, 3.75, 1.40, "LERF 2D+3D", "展示 rendered heatmap 和 direct 3D selection 是否一致。", COLORS["teal"], 14)
    d.card(s, 4.80, 1.35, 3.75, 1.40, "ScanNet point query", "展示开放词汇类别在点云中的空间定位。", COLORS["blue"], 14)
    d.card(s, 8.75, 1.35, 3.75, 1.40, "Failure / boundary", "展示 Waldo、小物体、边界碎片化等真实问题。", COLORS["amber"], 14)
    d.bullets(s, 1.00, 3.35, 11.20, 1.70, [
        "主文 qualitative 不宜堆太多，应围绕“同一 memory 支撑多协议 query”。",
        "补充材料可放更多 failure 和 support-policy 可视化，帮助解释 Direct3D caveat。",
        "导师需要决定：主文图优先强调强结果，还是强结果 + 典型失败各放一部分。",
    ], 15)
    d.note("定性页偏策略讨论，不强行塞所有图片。")

    # 19
    s = d.slide("当前边界和风险：要主动写清，不要等 reviewer 抓", "风险")
    d.card(s, 0.80, 1.28, 3.85, 1.48, "不是 universal feature superiority", "只说 selected primary downstream metrics；secondary caveats 单独列出。", COLORS["red"], 13)
    d.card(s, 4.75, 1.28, 3.85, 1.48, "不是 full ScanNet SOTA", "ScanNet 是 VALA-aligned 8-scene direct point-query protocol。", COLORS["amber"], 13)
    d.card(s, 8.70, 1.28, 3.85, 1.48, "不是外部后处理堆叠", "Direct3D compact row、support guard、SAM3 diagnostic 必须分层。", COLORS["blue"], 13)
    d.bullets(s, 1.05, 3.45, 11.2, 1.75, [
        "最危险的写法：只报最优数值，不解释 protocol / query interface / diagnostic 边界。",
        "最稳的写法：每条 claim 对应一个 main evidence，再把 caveat 放在同一节或 appendix。",
        "导师讨论重点：哪些 caveat 放主文，哪些放 supplement，避免削弱主贡献。",
    ], 15)
    d.note("这页是导师讨论版和对外短讲版最大的区别。")

    # 20
    s = d.slide("导师讨论点：投稿前需要拍板的 5 个问题", "下一步")
    d.table(s, 0.70, 1.20, 11.90, 4.60, [
        ["问题", "当前建议"],
        ["1. 主文是否放 frame-wise RADIO controlled table？", "建议放，能直接回应“是否只是压缩”的质疑。"],
        ["2. Direct3D 是否报告 SAM3-box diagnostic？", "可以放 supplement，主文避免把它当核心方法。"],
        ["3. ScanNet diagnostic 和 paper-facing row 如何并存？", "主文只放 paper-facing row，diagnostic 放 appendix/讨论。"],
        ["4. 消融表是否需要拆成 architecture / query interface / downstream？", "建议拆，避免一个大表读不出逻辑。"],
        ["5. 还需要补哪些实验？", "优先补能关闭 reviewer 质疑的控制实验，而不是无边界加指标。"],
    ], 10, emphasize_rows=[1, 2, 3, 4, 5])
    d.text(s, 1.00, 6.15, 11.20, 0.42, "最终目标：把“方法创新”和“证据闭环”同时讲清，而不是只展示一组好看的主结果。", 16, COLORS["teal"], bold=True, align=PP_ALIGN.CENTER)
    d.note("收尾用来开启导师讨论。")

    return d


def _build_deck_compact_previous() -> AdvisorDeck:
    d = AdvisorDeck()

    # 1
    s = d.slide("GaussFM：开放词汇三维场景的紧凑基础特征记忆", "定位")
    d.text(
        s,
        0.95,
        1.38,
        11.55,
        0.80,
        "核心主张：把 frame-wise RADIO 特征转化为可部署、可渲染、可直接查询的 3D Gaussian scene memory。",
        24,
        COLORS["ink"],
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    d.card(s, 0.95, 2.72, 3.55, 1.35, "研究对象", "不是新的 RGB 3DGS 重建器，而是基础特征场景记忆。", COLORS["red"], 14)
    d.card(s, 4.95, 2.72, 3.55, 1.35, "方法抓手", "低维 Gaussian code + 空间上下文 + 按需 RADIO-space 重建。", COLORS["teal"], 14)
    d.card(s, 8.95, 2.72, 3.55, 1.35, "验证闭环", "2D rendered query、3D primitive query、ScanNet point query。", COLORS["blue"], 14)
    d.text(s, 1.10, 4.85, 11.20, 0.48, "这版 PPT 的目标：短、直、完整，把方法创新和实验必要性讲成一条线。", 17, COLORS["muted"], align=PP_ALIGN.CENTER)
    d.note("开场只保留 thesis，不做过度包装。")

    # 2
    s = d.slide("问题：2D foundation features 到 3D memory 之间有断点", "问题")
    d.card(s, 0.85, 1.30, 3.72, 1.35, "视角局部性", "单帧 RADIO 特征受视角、遮挡和尺度影响，跨视角不天然一致。", COLORS["blue"], 14)
    d.card(s, 4.80, 1.30, 3.72, 1.35, "高维存储", "直接为每个 Gaussian 存 1280-D 特征，成本高且不优雅。", COLORS["amber"], 14)
    d.card(s, 8.75, 1.30, 3.72, 1.35, "查询割裂", "2D heatmap、3D primitive、point query 往往需要不同接口。", COLORS["red"], 14)
    d.text(s, 0.95, 3.35, 11.55, 0.62, "本文要回答的不是“能否压缩特征”，而是：能否学习一个紧凑的 3D 特征记忆，同时保留 foundation feature 的开放词汇可用性。", 20, COLORS["ink"], bold=True, align=PP_ALIGN.CENTER)
    d.bullets(s, 1.10, 4.55, 11.10, 1.25, [
        "如果只在每个 novel view 重新跑 2D encoder，方法就不是 3D memory。",
        "如果只能做 rendered-view query，方法就没有真正进入 primitive / point-level 3D 查询。",
        "如果没有原始 RADIO 和显式 1280-D memory 对照，compact claim 站不稳。",
    ], 14)
    d.note("这页把审稿人会问的三个断点提前摆出来。")

    # 3
    s = d.slide("核心思想与贡献：一个 compact memory，三种开放词汇接口", "贡献")
    d.card(s, 0.80, 1.25, 3.85, 2.05, "贡献 1", "Compact Gaussian Feature Memory：每个 Gaussian 存低维 code，按需重建 RADIO-compatible features。", COLORS["teal"], 14)
    d.card(s, 4.75, 1.25, 3.85, 2.05, "贡献 2", "多视角 primitive anchoring：把跨视角证据压回稳定的 3D support。", COLORS["blue"], 14)
    d.card(s, 8.70, 1.25, 3.85, 2.05, "贡献 3", "统一 query interface：同一场景记忆服务 view-space、primitive-space、point-space 查询。", COLORS["green"], 14)
    d.text(s, 1.00, 4.10, 11.35, 0.58, "讲故事的关键：每个方法组件都必须对应一个不可替代的问题，每个 claim 都必须有主结果、控制对比和核心消融支撑。", 19, COLORS["ink"], bold=True, align=PP_ALIGN.CENTER)
    d.chip(s, 2.15, 5.25, 2.05, "Compact storage", COLORS["teal"])
    d.chip(s, 4.65, 5.25, 2.05, "RADIO-space reconstruction", COLORS["blue"])
    d.chip(s, 7.15, 5.25, 2.05, "3D support", COLORS["green"])
    d.chip(s, 9.65, 5.25, 2.05, "Unified query", COLORS["violet"])
    d.note("贡献页同时承担叙事总纲。")

    # 4
    s = d.slide("方法总览：训练时重建 foundation features，推理时查询 compact memory", "方法")
    d.image(s, FRAMEWORK_FIG, 0.62, 1.16, 12.15, 4.85)
    d.text(s, 0.95, 6.14, 11.50, 0.42, "一句话：冻结 RGB/几何场景作为支架，学习共享 compact feature memory；推理时按 view / primitive / point 解码查询。", 15, COLORS["muted"], align=PP_ALIGN.CENTER)
    d.note("已替换为用户提供的新方法总框架图，并使用中性 GaussFM 标题。")

    # 5
    s = d.slide("方法核心：Compact Gaussian Feature Memory 的四个部件", "方法")
    d.image(s, MEMORY_FIG, 0.62, 1.14, 7.05, 4.98)
    d.card(s, 7.95, 1.24, 4.40, 0.80, "1. 低维 Gaussian code", "避免每个 primitive 存完整 1280-D 特征。", COLORS["teal"], 12)
    d.card(s, 7.95, 2.20, 4.40, 0.80, "2. 空间上下文", "用局部 3D context 补足单个 Gaussian 的局部性。", COLORS["blue"], 12)
    d.card(s, 7.95, 3.16, 4.40, 0.80, "3. 可靠性 / 可见性", "把噪声视角和不稳定 support 从 query interface 中剥离。", COLORS["green"], 12)
    d.card(s, 7.95, 4.12, 4.40, 0.80, "4. RADIO-space decoder", "在查询位置恢复 foundation-compatible feature。", COLORS["violet"], 12)
    d.text(s, 7.98, 5.35, 4.35, 0.56, "创新点不是“低维压缩”四个字，而是低维存储、空间聚合、可靠性建模和可查询解码的一体化。", 13, COLORS["ink"], bold=True)
    d.note("已替换为用户提供的新 compact memory 结构图。")

    # 6
    s = d.slide("实验设计：三条主结果，两类因果证据", "实验")
    d.table(s, 0.62, 1.20, 12.10, 4.75, [
        ["要证明的 claim", "主实验", "必须补上的证据"],
        ["view-space 查询可用", "LERF rendered-view 2D OVS", "完整 baseline 集 + 定性 heatmap"],
        ["primitive-space 查询可用", "LERF direct 3D object selection", "support calibration 消融"],
        ["point-space 查询可迁移", "VALA-aligned ScanNet point query", "协议边界与空间可视化"],
        ["不是最近帧 cache", "frame-wise RADIO / nearest-view cache", "same protocol controlled comparison"],
        ["不是粗暴 1280-D 存储", "per-Gaussian 1280-D RADIO memory", "storage 与性能同时报告"],
    ], 10, emphasize_rows=[1, 2, 3])
    d.text(s, 0.95, 6.20, 11.50, 0.36, "原则：主文/汇报不堆小消融，只保留能支撑核心贡献的强对照。", 15, COLORS["red"], bold=True, align=PP_ALIGN.CENTER)
    d.note("这页替代原先分散的实验导航。")

    # 7
    s = d.slide("主结果一：LERF rendered-view 2D OVS，完整 baseline 集下 mIoU 领先", "主结果")
    d.table(s, 0.72, 1.22, 7.10, 4.20, [
        ["Method", "mIoU", "Acc"],
        ["LangSplat", "51.40", "84.30"],
        ["GAGS", "54.12", "81.66"],
        ["OccamLGS", "61.30", "82.50"],
        ["GOI", "42.00", "59.20"],
        ["GALA", "55.49", "73.43"],
        ["LangSplatV2", "59.90", "84.10"],
        ["GaussFM", "64.98", "82.68"],
    ], 10, emphasize_rows=[7])
    d.card(s, 8.25, 1.35, 4.20, 1.20, "读法", "GaussFM 在 mean mIoU 上领先；Acc 与最强 baseline 接近，但不是最高。", COLORS["teal"], 13)
    d.card(s, 8.25, 2.95, 4.20, 1.20, "证明点", "compact memory 解码出的 dense feature map 可直接支持 novel-view 文本 grounding。", COLORS["blue"], 13)
    d.card(s, 8.25, 4.55, 4.20, 1.20, "写法边界", "不要只放三四个方法；这一页必须展示完整对比集合。", COLORS["amber"], 13)
    d.note("补全 LERF 2D OVS 对比方法。")

    # 8
    s = d.slide("主结果二：LERF direct 3D object selection，primitive-space 查询成立", "主结果")
    d.table(s, 0.78, 1.20, 6.85, 4.30, [
        ["Method", "mIoU", "Acc"],
        ["OpenGaussian", "38.36", "51.43"],
        ["SuperGSeg", "35.94", "52.02"],
        ["OccamLGS", "47.22", "74.84"],
        ["Dr. Splat", "43.29", "64.30"],
        ["GALA", "36.71", "59.71"],
        ["LangSplatV2", "35.87", "55.80"],
        ["GaussFM", "54.36", "80.84"],
    ], 10, emphasize_rows=[7])
    d.card(s, 8.10, 1.36, 4.35, 1.28, "证明点", "不是把 2D heatmap 反投影，而是在 Gaussian primitive 层面直接读出对象 support。", COLORS["teal"], 13)
    d.card(s, 8.10, 3.06, 4.35, 1.28, "必须解释", "Direct3D 的 query interface calibration 是核心支持策略，需要用消融单独闭环。", COLORS["red"], 13)
    d.note("补全 LERF direct 3D 对比方法。")

    # 9
    s = d.slide("主结果三：VALA-aligned ScanNet point query，跨数据集仍有效", "主结果")
    d.table(s, 0.55, 1.16, 12.25, 4.40, [
        ["Method", "19 mIoU / mAcc", "15 mIoU / mAcc", "10 mIoU / mAcc"],
        ["LangSplat", "2.45 / 8.59", "3.45 / 13.21", "6.48 / 21.89"],
        ["LangSplatV2", "14.75 / 25.47", "17.09 / 35.68", "22.83 / 41.52"],
        ["OpenGaussian", "27.73 / 42.01", "29.67 / 46.15", "39.93 / 57.34"],
        ["Dr. Splat", "29.31 / 47.68", "33.25 / 54.33", "44.19 / 65.19"],
        ["OccamLGS", "31.93 / 48.93", "34.25 / 53.71", "45.16 / 64.39"],
        ["VALA", "32.11 / 50.05", "35.10 / 54.77", "46.21 / 65.61"],
        ["GaussFM", "36.55 / 50.57", "42.78 / 72.85", "57.85 / 77.93"],
    ], 8, emphasize_rows=[7])
    d.text(s, 0.85, 6.00, 11.75, 0.44, "写法必须保守：这是 VALA-aligned 8-scene direct point-query protocol，不是完整 ScanNet semantic segmentation leaderboard。", 15, COLORS["red"], bold=True, align=PP_ALIGN.CENTER)
    d.note("补全 ScanNet 对比方法，并明确协议边界。")

    # 10
    s = d.slide("定性可视化：同一 memory 输出 2D、3D 和 point-level 结果", "定性")
    d.image(s, QUAL_LERF, 0.55, 1.22, 6.20, 4.55)
    d.image(s, QUAL_SCAN, 6.95, 1.22, 5.85, 4.55)
    d.text(s, 0.80, 5.95, 5.75, 0.38, "LERF：rendered heatmap 与 direct 3D selection 指向同一目标。", 12, COLORS["muted"], align=PP_ALIGN.CENTER)
    d.text(s, 7.15, 5.95, 5.45, 0.38, "ScanNet：point query 展示开放词汇类别的空间定位。", 12, COLORS["muted"], align=PP_ALIGN.CENTER)
    d.note("补上实际 qualitative visualization，而不是只讲策略。")

    # 11
    s = d.slide("原始 RADIO 控制对比：优势不是最近帧 cache，也不是显式 1280-D 存储", "控制实验")
    d.table(s, 0.55, 1.18, 8.55, 4.38, [
        ["Method", "Compact", "3D memory", "LERF LocAcc", "LERF mIoU", "Storage"],
        ["Frame-wise RADIO", "no", "no", "0.7985", "0.4634", "per-frame cache"],
        ["Nearest-view RADIO cache", "no", "no", "0.2722", "0.1545", "per-frame cache"],
        ["Per-Gaussian 1280-D RADIO", "no", "yes", "0.5642", "0.3182", "1039.7 MiB mean"],
        ["GaussFM core", "yes", "yes", "0.8598", "0.5707", "compact checkpoint"],
        ["GaussFM + boundary", "yes", "yes", "0.8598", "0.5889", "same memory + query interface"],
    ], 7, emphasize_rows=[4, 5])
    d.card(s, 9.45, 1.45, 3.00, 1.30, "直接回答", "compact field 是否只是低质压缩？否。", COLORS["teal"], 14)
    d.card(s, 9.45, 3.05, 3.00, 1.30, "直接回答", "是否靠最近训练帧投机？否。", COLORS["blue"], 14)
    d.card(s, 9.45, 4.65, 3.00, 1.30, "直接回答", "是否需要存原始 1280-D？否。", COLORS["green"], 14)
    d.note("保留原始 RADIO 对比，避免 compact claim 被质疑。")

    # 12
    s = d.slide("冻结头下游 probes：证明 feature 可用性，但不夸大为 universal superiority", "控制实验")
    d.table(s, 0.50, 1.14, 12.35, 4.45, [
        ["Task", "Primary", "Frame-wise RADIO", "GaussFM", "Delta", "Secondary caveat"],
        ["SAM3 point prompt", "mIoU", "0.3700", "0.4173", "+0.0473", "LocAcc same"],
        ["SAM3 box prompt", "mIoU", "0.6560", "0.6638", "+0.0078", "LocAcc lower"],
        ["SAM3 mask propagation", "mIoU", "0.3583", "0.3756", "+0.0173", "LocAcc lower"],
        ["DINOv3 dense matching", "Mean score", "0.8547", "0.9048", "+0.0501", "HitRate lower"],
        ["DINOv3 mask propagation", "mIoU", "0.4606", "0.4677", "+0.0071", "LocAcc higher"],
    ], 8, emphasize_rows=[1, 2, 3, 4, 5])
    d.text(s, 0.80, 5.88, 11.80, 0.46, "Claim-safe wording：selected primary feature-usability metrics improve；secondary caveats 同页呈现，不宣称无条件全胜。", 15, COLORS["red"], bold=True, align=PP_ALIGN.CENTER)
    d.note("保留下游 probes，但降格为 feature usability evidence。")

    # 13
    s = d.slide("核心消融：只围绕三项突出贡献，而不是罗列小模块", "消融")
    d.table(s, 0.52, 1.18, 12.35, 4.58, [
        ["核心贡献", "弱化 / 移除设置", "主设置", "关键变化", "解释"],
        ["Foundation-space reconstruction", "w/o HCD / CTR", "Full GaussFM", "LocAcc 0.531 -> 0.858; mIoU 0.260 -> 0.485", "compact code 必须能回到 RADIO space"],
        ["Multi-view primitive anchoring", "w/o VPR-to-field registration", "registered field", "mIoU 0.046 -> 0.412; Acc 0.071 -> 0.588", "3D support 不是自然涌现"],
        ["Support-calibrated query interface", "pure compact prompt ensemble", "score-component guard", "mIoU 0.457 -> 0.501; Acc 0.685 -> 0.704", "Direct3D 需要稳定 support 策略"],
    ], 7, emphasize_rows=[1, 2, 3])
    d.text(s, 0.90, 6.05, 11.60, 0.42, "FDH、refiner、hybrid、ScanNet diagnostic 等小点放 supplement 或备注；主线消融只证明核心创新不可替代。", 15, COLORS["muted"], align=PP_ALIGN.CENTER)
    d.note("回应用户反馈：消融从碎片小点收敛到核心贡献。")

    # 14
    s = d.slide("Storage / efficiency：同时报告 latent、package、full checkpoint", "效率")
    d.table(s, 0.55, 1.16, 12.20, 4.15, [
        ["Scene", "Direct 1280-D", "Latent payload", "Latent saving", "Feature package", "Package saving", "Full saving"],
        ["Figurines", "412.1 MiB", "20.6 MiB", "20.00x", "199.0 MiB", "2.07x", "1.74x"],
        ["Ramen", "934.3 MiB", "46.7 MiB", "20.00x", "225.1 MiB", "4.15x", "3.00x"],
        ["Teatime", "1123.4 MiB", "56.2 MiB", "20.00x", "234.5 MiB", "4.79x", "3.32x"],
        ["Waldo Kitchen", "1688.8 MiB", "84.4 MiB", "20.00x", "262.8 MiB", "6.43x", "4.04x"],
    ], 8, emphasize_cols=[3])
    d.card(s, 0.95, 5.68, 5.65, 0.82, "汇报口径", "latent payload 是最干净的 feature-memory 压缩；package/full checkpoint 是部署视角。", COLORS["teal"], 13)
    d.card(s, 6.95, 5.68, 5.25, 0.82, "风险", "只报 full checkpoint 会低估 compact feature-memory 的真实优势。", COLORS["amber"], 13)
    d.note("这页用于和导师讨论论文表格该报哪几列。")

    # 15
    s = d.slide("当前边界和写法：主动收束 claim，避免被 reviewer 抓漏洞", "边界")
    d.card(s, 0.80, 1.28, 3.85, 1.48, "不是 universal feature superiority", "只说 selected primary downstream metrics；secondary caveats 单独列出。", COLORS["red"], 13)
    d.card(s, 4.75, 1.28, 3.85, 1.48, "不是 full ScanNet SOTA", "ScanNet 是 VALA-aligned 8-scene direct point-query protocol。", COLORS["amber"], 13)
    d.card(s, 8.70, 1.28, 3.85, 1.48, "不是外部后处理堆叠", "Direct3D compact row、support guard、diagnostic 必须分层。", COLORS["blue"], 13)
    d.bullets(s, 1.05, 3.45, 11.2, 1.75, [
        "最危险写法：只报最优数值，不解释 protocol / query interface / diagnostic 边界。",
        "最稳写法：每条 claim 对应一个 main evidence，再把 caveat 放在同一节或 appendix。",
        "命名口径要统一：投稿材料只使用 GaussFM 与 compact Gaussian feature memory。",
    ], 15)
    d.note("这页是导师讨论版保留的审稿风险页。")

    # 16
    s = d.slide("结论页：故事线、证据线、投稿前动作", "下一步")
    d.card(s, 0.78, 1.30, 3.85, 2.00, "故事线", "2D foundation features 不能直接成为 3D memory；GaussFM 学习紧凑、可重建、可查询的 Gaussian feature memory。", COLORS["teal"], 13)
    d.card(s, 4.73, 1.30, 3.85, 2.00, "证据线", "三张主结果表补全 baseline；定性可视化展示多接口；RADIO 控制对比和核心消融关闭主要质疑。", COLORS["blue"], 13)
    d.card(s, 8.68, 1.30, 3.85, 2.00, "投稿前动作", "主文保留三项核心消融；小型诊断、完整 query interface registry 和更多 failure cases 放 supplement。", COLORS["green"], 13)
    d.text(s, 0.95, 4.45, 11.50, 0.72, "最终目标：让大小同行先理解“为什么需要 3D foundation-feature memory”，再相信“GaussFM 的 compact memory 确实可用且必要”。", 20, COLORS["ink"], bold=True, align=PP_ALIGN.CENTER)
    d.note("收尾回到完整故事，而不是停在实验列表。")

    return d


def build_deck() -> AdvisorDeck:
    d = AdvisorDeck()

    # 1
    s = d.slide("GaussFM：开放词汇三维场景的紧凑基础特征记忆", "Title")
    d.text(
        s,
        0.90,
        1.36,
        11.60,
        0.82,
        "把 frame-wise RADIO 特征转化为可部署、可重建、可直接查询的 3D Gaussian scene memory。",
        24,
        COLORS["ink"],
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    d.card(s, 0.98, 2.82, 3.45, 1.25, "Representation", "compact Gaussian feature memory", COLORS["teal"], 15)
    d.card(s, 4.95, 2.82, 3.45, 1.25, "Training signal", "frozen RADIO feature reconstruction", COLORS["blue"], 15)
    d.card(s, 8.92, 2.82, 3.45, 1.25, "Tasks", "2D / 3D / point-level open-vocabulary query", COLORS["green"], 15)
    d.text(s, 1.10, 4.90, 11.10, 0.44, "会议报告逻辑：背景问题 -> 现有方法缺口 -> 方法设计 -> 实验证据 -> 边界与结论。", 16, COLORS["muted"], align=PP_ALIGN.CENTER)
    d.note("首页给出论文 thesis 和报告路线。")

    # 2
    s = d.slide("任务背景：foundation features 正在成为 3D 场景理解的通用接口", "Background")
    d.card(s, 0.85, 1.28, 3.70, 1.48, "2D foundation models", "RADIO / DINO / SAM / SigLIP 等提供强 dense features 与冻结头能力。", COLORS["blue"], 14)
    d.card(s, 4.82, 1.28, 3.70, 1.48, "3D Gaussian scenes", "显式几何和高效渲染适合做可部署的场景表示。", COLORS["teal"], 14)
    d.card(s, 8.80, 1.28, 3.70, 1.48, "Open vocabulary queries", "用户希望用文本或类别在 novel view、primitive 和 point level 查询目标。", COLORS["green"], 14)
    d.text(s, 0.95, 3.42, 11.55, 0.64, "核心机会：如果 3D 场景本身能携带 foundation-compatible features，就不必每次查询都回到逐帧 2D 特征。", 20, COLORS["ink"], bold=True, align=PP_ALIGN.CENTER)
    d.bullets(s, 1.10, 4.55, 11.10, 1.16, [
        "这不是 RGB 重建问题，而是 3D scene memory 的表示学习问题。",
        "评价也不能只看 novel-view heatmap，还要看 primitive-level 和 point-level 查询。",
    ], 15)
    d.note("标准会议报告先解释任务和研究对象。")

    # 3
    s = d.slide("动机：逐帧 2D 特征强，但不是可部署的 3D 记忆", "Motivation")
    d.card(s, 0.82, 1.22, 3.75, 1.55, "View-local", "每帧特征受视角、遮挡、尺度影响；跨视角一致性不是天然属性。", COLORS["blue"], 14)
    d.card(s, 4.78, 1.22, 3.75, 1.55, "Storage-heavy", "直接存 per-view 或 per-Gaussian 1280-D 特征，成本随场景规模快速上升。", COLORS["amber"], 14)
    d.card(s, 8.74, 1.22, 3.75, 1.55, "Interface mismatch", "2D heatmap 难以直接回答 3D primitive / point query。", COLORS["red"], 14)
    d.text(s, 0.95, 3.55, 11.50, 0.66, "问题转化：我们需要的不是一个更强的 2D encoder，而是一个能把多视角 foundation evidence 压进 3D 的紧凑场景记忆。", 20, COLORS["ink"], bold=True, align=PP_ALIGN.CENTER)
    d.note("动机页明确为什么不能停留在 frame-wise feature extraction。")

    # 4
    s = d.slide("相关方法脉络：已有工作各解决了一部分，但接口仍然割裂", "Related Work")
    d.table(s, 0.66, 1.24, 12.00, 4.70, [
        ["路线", "代表能力", "仍然缺什么"],
        ["Rendered language fields", "novel-view 2D text grounding", "通常不直接提供 primitive / point-level 3D query"],
        ["Open-vocabulary 3DGS", "把语义或语言信息绑定到 Gaussian / object", "高维存储、跨视角一致性和查询接口常分开处理"],
        ["Instance / proposal methods", "对象级支持更稳定", "容易依赖后处理或任务特定 proposal"],
        ["Frame-wise foundation features", "冻结模型能力强、迁移性好", "不是压缩后的 3D scene memory"],
    ], 11, emphasize_rows=[4])
    d.text(s, 0.95, 6.10, 11.45, 0.40, "GaussFM 的定位：学习一个 compact foundation-feature Gaussian memory，而不是再做一个任务特定分类器或 per-view feature cache。", 15, COLORS["teal"], bold=True, align=PP_ALIGN.CENTER)
    d.note("补上标准报告中的 related work overview。")

    # 5
    s = d.slide("问题分析：顶会论文必须证明三个层面的“不是”", "Gap")
    d.card(s, 0.85, 1.30, 3.70, 1.55, "不是 per-view cache", "需要证明优于 frame-wise / nearest-view RADIO，而非简单查最近训练帧。", COLORS["red"], 14)
    d.card(s, 4.82, 1.30, 3.70, 1.55, "不是粗暴高维存储", "需要证明 compact memory 比 per-Gaussian 1280-D 特征更有效、更轻。", COLORS["amber"], 14)
    d.card(s, 8.80, 1.30, 3.70, 1.55, "不是单一任务技巧", "需要同时覆盖 view-space、primitive-space、point-space 查询。", COLORS["blue"], 14)
    d.text(s, 1.00, 4.05, 11.30, 0.62, "因此实验不能只是一张好看的主表；必须包含完整 baseline、原始 RADIO 控制对比、核心组件消融和定性证据。", 20, COLORS["ink"], bold=True, align=PP_ALIGN.CENTER)
    d.note("这页解释为什么后续实验设计是必要的。")

    # 6
    s = d.slide("本文贡献：一个记忆，三个接口，两类闭环证据", "Contribution")
    d.card(s, 0.82, 1.24, 3.75, 2.05, "C1. 紧凑记忆", "每个 Gaussian 存低维 code，并结合空间上下文、质量和可见性信息。", COLORS["teal"], 14)
    d.card(s, 4.78, 1.24, 3.75, 2.05, "C2. 特征重建", "按需重建 RADIO-compatible features，使冻结下游能力可以继续使用。", COLORS["blue"], 14)
    d.card(s, 8.74, 1.24, 3.75, 2.05, "C3. 多粒度查询", "同一场景记忆支持 rendered-view、primitive-level、point-level open-vocabulary query。", COLORS["green"], 14)
    d.table(s, 1.35, 4.18, 10.60, 1.25, [
        ["证据类型", "回答的问题"],
        ["主结果", "是否在完整 baseline 下有效？"],
        ["控制对比 + 消融", "为什么不是 cache / 高维存储 / 任务技巧？"],
    ], 12, emphasize_rows=[1, 2])
    d.note("贡献页把方法贡献和证据贡献绑定。")

    # 7
    s = d.slide("方法总览：从 posed RGB views 到 compact Gaussian feature memory", "Method")
    d.image(s, FRAMEWORK_FIG, 0.60, 1.14, 12.15, 4.86)
    d.text(s, 0.92, 6.12, 11.55, 0.42, "训练：重建冻结 RADIO feature 并注册多视角 primitive evidence；推理：从同一 compact memory 解码到不同查询位置。", 15, COLORS["muted"], align=PP_ALIGN.CENTER)
    d.note("使用新总框架图。")

    # 8
    s = d.slide("方法核心：Compact Gaussian Feature Memory", "Method")
    d.image(s, MEMORY_FIG, 0.58, 1.12, 7.10, 5.02)
    d.card(s, 7.92, 1.22, 4.42, 0.82, "低维 Gaussian codes", "每个 primitive 存 compact latent，而非完整 1280-D feature。", COLORS["teal"], 12)
    d.card(s, 7.92, 2.22, 4.42, 0.82, "Hybrid 3D context", "结合 view-space detail 与 spatial 3D context。", COLORS["blue"], 12)
    d.card(s, 7.92, 3.22, 4.42, 0.82, "Reliability / visibility", "显式建模质量与可见性，稳定 3D 支持区域。", COLORS["green"], 12)
    d.card(s, 7.92, 4.22, 4.42, 0.82, "RADIO-space reconstruction", "查询时恢复 foundation-compatible feature。", COLORS["violet"], 12)
    d.text(s, 7.98, 5.42, 4.30, 0.50, "关键点：compact 不是简单降维，而是可查询的场景级 foundation-feature memory。", 13, COLORS["ink"], bold=True)
    d.note("使用新 compact Gaussian feature memory 图。")

    # 9
    s = d.slide("查询接口：同一个 memory 服务三种评测协议", "Method")
    d.table(s, 0.65, 1.22, 12.05, 4.55, [
        ["查询空间", "输入", "输出", "对应实验"],
        ["View-space", "novel-view camera + text", "2D heatmap / mask", "LERF rendered-view 2D OVS"],
        ["Primitive-space", "text query over Gaussian primitives", "stable object support", "LERF direct 3D object selection"],
        ["Point-space", "3D point / point cloud query", "class score / binary support", "VALA-aligned ScanNet point query"],
    ], 12, emphasize_rows=[1, 2, 3])
    d.text(s, 0.90, 6.00, 11.60, 0.46, "用词边界：这里称为“查询接口/查询解码/支持校准”，不把它包装成额外的下游任务头。", 15, COLORS["red"], bold=True, align=PP_ALIGN.CENTER)
    d.note("替代不严谨的 query interface 表述。")

    # 10
    s = d.slide("实验设置：主结果、控制对比和消融各回答一个问题", "Experiments")
    d.table(s, 0.62, 1.20, 12.10, 4.72, [
        ["实验块", "协议", "它回答什么"],
        ["LERF rendered-view 2D OVS", "same-protocol reproduced baselines", "novel-view 文本定位是否有效"],
        ["LERF direct 3D selection", "primitive query-select-render protocol", "是否能直接在 3D support 上查询"],
        ["ScanNet point query", "VALA-aligned 8-scene protocol", "是否迁移到点级开放词汇查询"],
        ["RADIO controls", "frame-wise / nearest cache / 1280-D memory", "是否只是 cache 或高维存储"],
        ["Core ablations", "reconstruction / registration / support calibration", "核心组件是否必要"],
    ], 10, emphasize_rows=[1, 2, 3])
    d.note("实验设置页让后续结果有导航。")

    # 11
    s = d.slide("主结果一：LERF rendered-view 2D OVS", "Results")
    d.table(s, 0.90, 1.20, 6.85, 4.45, [
        ["Method", "mIoU", "Acc"],
        ["LangSplat", "51.40", "84.30"],
        ["GAGS", "54.12", "81.66"],
        ["OccamLGS", "61.30", "82.50"],
        ["GOI", "42.00", "59.20"],
        ["GALA", "55.49", "73.43"],
        ["LangSplatV2", "59.90", "84.10"],
        ["GaussFM", "64.98", "82.68"],
    ], 10, emphasize_rows=[7])
    d.card(s, 8.25, 1.42, 4.10, 1.20, "结论", "完整 baseline 集下，GaussFM 的 mean mIoU 最高。", COLORS["teal"], 14)
    d.card(s, 8.25, 3.10, 4.10, 1.20, "读法", "Acc 接近最强方法但不是最高，因此主张应聚焦定位质量而非所有指标全胜。", COLORS["amber"], 14)
    d.note("表格字号提高到 10。")

    # 12
    s = d.slide("主结果二：LERF direct 3D object selection", "Results")
    d.table(s, 0.90, 1.20, 6.85, 4.45, [
        ["Method", "mIoU", "Acc"],
        ["OpenGaussian", "38.36", "51.43"],
        ["SuperGSeg", "35.94", "52.02"],
        ["OccamLGS", "47.22", "74.84"],
        ["Dr. Splat", "43.29", "64.30"],
        ["GALA", "36.71", "59.71"],
        ["LangSplatV2", "35.87", "55.80"],
        ["GaussFM", "54.36", "80.84"],
    ], 10, emphasize_rows=[7])
    d.card(s, 8.25, 1.40, 4.10, 1.25, "结论", "同一 memory 可以在 primitive-space 形成稳定对象支持。", COLORS["teal"], 14)
    d.card(s, 8.25, 3.15, 4.10, 1.25, "意义", "这与 rendered-view 2D heatmap 是不同能力，证明方法不只停在图像空间。", COLORS["blue"], 14)
    d.note("Direct 3D 表补全对比方法。")

    # 13
    s = d.slide("主结果三：VALA-aligned ScanNet point query", "Results")
    d.table(s, 0.50, 1.14, 12.30, 4.55, [
        ["Method", "19 mIoU / mAcc", "15 mIoU / mAcc", "10 mIoU / mAcc"],
        ["LangSplat", "2.45 / 8.59", "3.45 / 13.21", "6.48 / 21.89"],
        ["LangSplatV2", "14.75 / 25.47", "17.09 / 35.68", "22.83 / 41.52"],
        ["OpenGaussian", "27.73 / 42.01", "29.67 / 46.15", "39.93 / 57.34"],
        ["Dr. Splat", "29.31 / 47.68", "33.25 / 54.33", "44.19 / 65.19"],
        ["OccamLGS", "31.93 / 48.93", "34.25 / 53.71", "45.16 / 64.39"],
        ["VALA", "32.11 / 50.05", "35.10 / 54.77", "46.21 / 65.61"],
        ["GaussFM", "36.55 / 50.57", "42.78 / 72.85", "57.85 / 77.93"],
    ], 9, emphasize_rows=[7])
    d.text(s, 0.82, 6.00, 11.75, 0.44, "边界：这是 VALA-aligned 8-scene direct point-query protocol，不写成 full ScanNet semantic segmentation leaderboard。", 15, COLORS["red"], bold=True, align=PP_ALIGN.CENTER)
    d.note("ScanNet 表字号提升到 9，并保留完整方法集合。")

    # 14
    s = d.slide("定性一：同一 memory 同时支持 2D 与 3D 开放词汇查询", "Qualitative")
    d.image(s, QUAL_LERF, 0.55, 1.14, 12.15, 5.25)
    d.text(s, 0.95, 6.42, 11.35, 0.30, "LERF qualitative：rendered-view heatmap 与 direct 3D selection 在同一目标上形成一致支持。", 13, COLORS["muted"], align=PP_ALIGN.CENTER)
    d.note("补足主结果对应定性图。")

    # 15
    s = d.slide("定性二：GaussFM 与 frame-wise RADIO 的 visual comparison", "Qualitative")
    d.image(s, QUAL_RADIO, 0.62, 1.12, 12.00, 5.35)
    d.text(s, 0.95, 6.42, 11.35, 0.30, "该图直接服务控制问题：重建场景特征是否仍保留 frame-wise RADIO 的开放词汇定位能力。", 13, COLORS["muted"], align=PP_ALIGN.CENTER)
    d.note("加入原始 RADIO 定性对比。")

    # 16
    s = d.slide("定性三：冻结头 probes 中的 scene-feature usability", "Qualitative")
    d.image(s, QUAL_SAM_DINO, 0.55, 1.10, 12.20, 5.36)
    d.text(s, 0.95, 6.42, 11.35, 0.30, "冻结 SAM/DINO/SigLIP probes 进一步说明：GaussFM field 不是只为单个文本定位指标过拟合。", 13, COLORS["muted"], align=PP_ALIGN.CENTER)
    d.note("加入 SAM/DINO 原始 RADIO 对比定性图。")

    # 17
    s = d.slide("控制对比：不是最近帧 cache，也不是显式 1280-D Gaussian memory", "Controls")
    d.table(s, 0.58, 1.18, 8.70, 4.38, [
        ["Method", "Compact", "3D memory", "LocAcc", "mIoU"],
        ["Frame-wise RADIO", "no", "no", "0.7985", "0.4634"],
        ["Nearest-view RADIO cache", "no", "no", "0.2722", "0.1545"],
        ["Per-Gaussian 1280-D RADIO", "no", "yes", "0.5642", "0.3182"],
        ["GaussFM core", "yes", "yes", "0.8598", "0.5707"],
        ["GaussFM + boundary calibration", "yes", "yes", "0.8598", "0.5889"],
    ], 9, emphasize_rows=[4, 5])
    d.card(s, 9.55, 1.45, 2.95, 1.30, "结论 1", "不是最近帧 cache。", COLORS["blue"], 14)
    d.card(s, 9.55, 3.05, 2.95, 1.30, "结论 2", "不是粗暴存 1280-D。", COLORS["teal"], 14)
    d.note("控制对比表字号从 7 提升到 9。")

    # 18
    s = d.slide("核心消融：只保留支撑主要贡献的三项强对照", "Ablation")
    d.table(s, 0.55, 1.22, 12.20, 4.45, [
        ["核心贡献", "弱化设置", "主设置", "关键变化"],
        ["Foundation-space reconstruction", "w/o HCD / CTR", "Full GaussFM", "LocAcc 0.531 -> 0.858; mIoU 0.260 -> 0.485"],
        ["Multiview primitive registration", "w/o registration", "registered field", "mIoU 0.046 -> 0.412; Acc 0.071 -> 0.588"],
        ["Support-calibrated primitive query", "pure compact prompt ensemble", "score-component calibration", "mIoU 0.457 -> 0.501; Acc 0.685 -> 0.704"],
    ], 9, emphasize_rows=[1, 2, 3])
    d.text(s, 0.90, 6.05, 11.55, 0.42, "FDH、refiner、hybrid、ScanNet diagnostic 等小点放 supplement；主讲只证明核心创新不可替代。", 15, COLORS["muted"], align=PP_ALIGN.CENTER)
    d.note("消融表字号提升，内容收束为核心贡献。")

    # 19
    s = d.slide("Storage / efficiency：compact memory 的实际收益", "Efficiency")
    d.table(s, 0.55, 1.16, 12.20, 4.15, [
        ["Scene", "Direct 1280-D", "Latent payload", "Latent saving", "Feature package", "Full saving"],
        ["Figurines", "412.1 MiB", "20.6 MiB", "20.00x", "199.0 MiB", "1.74x"],
        ["Ramen", "934.3 MiB", "46.7 MiB", "20.00x", "225.1 MiB", "3.00x"],
        ["Teatime", "1123.4 MiB", "56.2 MiB", "20.00x", "234.5 MiB", "3.32x"],
        ["Waldo Kitchen", "1688.8 MiB", "84.4 MiB", "20.00x", "262.8 MiB", "4.04x"],
    ], 9, emphasize_cols=[3])
    d.card(s, 0.95, 5.68, 5.65, 0.82, "汇报口径", "latent payload 体现 feature-memory 压缩；package/full checkpoint 体现部署视角。", COLORS["teal"], 13)
    d.card(s, 6.95, 5.68, 5.25, 0.82, "主张", "效率收益必须和性能控制对比一起讲，否则会被看成普通压缩。", COLORS["amber"], 13)
    d.note("效率表删掉一列以提升字号。")

    # 20
    s = d.slide("边界：哪些 claim 必须主动收束", "Discussion")
    d.card(s, 0.80, 1.28, 3.85, 1.48, "不是 universal feature superiority", "只说 selected primary downstream metrics；secondary caveats 单独列出。", COLORS["red"], 13)
    d.card(s, 4.75, 1.28, 3.85, 1.48, "不是 full ScanNet SOTA", "ScanNet 是 VALA-aligned direct point-query protocol。", COLORS["amber"], 13)
    d.card(s, 8.70, 1.28, 3.85, 1.48, "不是额外任务头堆叠", "Direct3D compact query、support calibration、diagnostic rows 必须分层。", COLORS["blue"], 13)
    d.bullets(s, 1.05, 3.45, 11.20, 1.75, [
        "最危险写法：只报最优数值，不解释 protocol / query interface / diagnostic 边界。",
        "最稳写法：每条 claim 对应一个 main evidence，再把 caveat 放在同一节或 appendix。",
        "术语口径：使用 GaussFM、compact Gaussian feature memory、query interface、support calibration。",
    ], 15)
    d.note("讨论页避免过度声明。")

    # 21
    s = d.slide("Takeaway：从 frame-wise features 到 queryable 3D scene memory", "Conclusion")
    d.card(s, 0.82, 1.35, 3.75, 1.90, "Problem", "2D foundation features 强，但 view-local、高维、不能直接构成 3D memory。", COLORS["red"], 14)
    d.card(s, 4.78, 1.35, 3.75, 1.90, "Method", "GaussFM 学习 compact Gaussian feature memory，并按需重建 RADIO-compatible features。", COLORS["teal"], 14)
    d.card(s, 8.74, 1.35, 3.75, 1.90, "Evidence", "完整 baseline、定性、RADIO controls、核心消融共同支撑三种查询接口。", COLORS["blue"], 14)
    d.text(s, 0.95, 4.55, 11.45, 0.72, "一句话收尾：GaussFM 把 3D Gaussian scene 从几何/外观表示推进为紧凑、可复用的 foundation-feature scene memory。", 20, COLORS["ink"], bold=True, align=PP_ALIGN.CENTER)
    d.note("会议报告式收尾。")

    return d


def main() -> None:
    deck = build_deck()
    deck.save()
    print(f"Wrote {OUT_PPTX}")
    print(f"Wrote {OUT_OUTLINE}")


if __name__ == "__main__":
    main()
