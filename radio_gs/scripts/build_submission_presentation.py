#!/usr/bin/env python3
"""Build a paper-submission presentation deck for RADIO-GS."""

from __future__ import annotations

import json
from pathlib import Path
from typing import Iterable, Sequence

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE, MSO_VERTICAL_ANCHOR
from pptx.util import Inches, Pt


REPO_ROOT = Path(__file__).resolve().parents[2]
PAPER_DIR = REPO_ROOT / "paper"
OUT_PPTX = PAPER_DIR / "radio_gs_submission_presentation.pptx"
OUT_OUTLINE = PAPER_DIR / "radio_gs_submission_presentation_outline.md"

SLIDE_W = 13.333
SLIDE_H = 7.5

COLORS = {
    "ink": RGBColor(24, 31, 42),
    "muted": RGBColor(89, 103, 118),
    "paper": RGBColor(248, 250, 252),
    "panel": RGBColor(255, 255, 255),
    "line": RGBColor(213, 219, 226),
    "blue": RGBColor(44, 93, 165),
    "teal": RGBColor(20, 133, 132),
    "amber": RGBColor(208, 134, 44),
    "red": RGBColor(186, 72, 72),
    "green": RGBColor(75, 143, 94),
    "violet": RGBColor(102, 88, 170),
}

FONT = "Microsoft YaHei"
FONT_LATIN = "Aptos"


class DeckBuilder:
    def __init__(self) -> None:
        self.prs = Presentation()
        self.prs.slide_width = Inches(SLIDE_W)
        self.prs.slide_height = Inches(SLIDE_H)
        self.blank = self.prs.slide_layouts[6]
        self.outline: list[tuple[str, list[str]]] = []

    def add_slide(self, title: str, section: str | None = None) -> object:
        slide = self.prs.slides.add_slide(self.blank)
        bg = slide.background
        bg.fill.solid()
        bg.fill.fore_color.rgb = COLORS["paper"]
        if section:
            self.text(slide, 0.45, 0.18, 3.2, 0.25, section.upper(), 8, COLORS["teal"], bold=True)
        self.text(slide, 0.45, 0.45, 9.9, 0.45, title, 22, COLORS["ink"], bold=True)
        self.line(slide, 0.45, 1.02, 12.45, 1.02, COLORS["line"])
        self.footer(slide)
        self.outline.append((title, []))
        return slide

    def add_note(self, text: str) -> None:
        if self.outline:
            self.outline[-1][1].append(text)

    def footer(self, slide: object) -> None:
        idx = len(self.prs.slides)
        self.text(slide, 0.45, 7.08, 2.4, 0.18, "RADIO-GS / Submission Deck", 7, COLORS["muted"])
        self.text(slide, 12.05, 7.08, 0.7, 0.18, str(idx), 7, COLORS["muted"], align=PP_ALIGN.RIGHT)

    def text(
        self,
        slide: object,
        x: float,
        y: float,
        w: float,
        h: float,
        value: str,
        size: int,
        color: RGBColor | None = None,
        bold: bool = False,
        align: int | None = None,
    ) -> object:
        box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = box.text_frame
        tf.clear()
        tf.word_wrap = True
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        p = tf.paragraphs[0]
        if align is not None:
            p.alignment = align
        r = p.add_run()
        r.text = value
        r.font.name = FONT
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color or COLORS["ink"]
        return box

    def bullets(
        self,
        slide: object,
        x: float,
        y: float,
        w: float,
        h: float,
        items: Sequence[str],
        size: int = 15,
        color: RGBColor | None = None,
        gap: float = 0.07,
    ) -> None:
        box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = box.text_frame
        tf.clear()
        tf.word_wrap = True
        tf.margin_left = Inches(0.05)
        tf.margin_right = Inches(0.05)
        tf.margin_top = Inches(0.02)
        tf.margin_bottom = Inches(0.02)
        for i, item in enumerate(items):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = item
            p.level = 0
            p.space_after = Pt(6 + gap * 10)
            p.font.name = FONT
            p.font.size = Pt(size)
            p.font.color.rgb = color or COLORS["ink"]

    def card(
        self,
        slide: object,
        x: float,
        y: float,
        w: float,
        h: float,
        title: str,
        body: str,
        accent: RGBColor,
        title_size: int = 12,
        body_size: int = 18,
    ) -> None:
        shape = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            Inches(x),
            Inches(y),
            Inches(w),
            Inches(h),
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = COLORS["panel"]
        shape.line.color.rgb = COLORS["line"]
        shape.adjustments[0] = 0.08
        accent_bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x), Inches(y), Inches(0.08), Inches(h))
        accent_bar.fill.solid()
        accent_bar.fill.fore_color.rgb = accent
        accent_bar.line.fill.background()
        self.text(slide, x + 0.22, y + 0.18, w - 0.35, 0.25, title, title_size, COLORS["muted"], bold=True)
        self.text(slide, x + 0.22, y + 0.55, w - 0.35, h - 0.65, body, body_size, COLORS["ink"], bold=True)

    def line(self, slide: object, x1: float, y1: float, x2: float, y2: float, color: RGBColor) -> None:
        line = slide.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
        line.line.color.rgb = color
        line.line.width = Pt(1.0)

    def table(
        self,
        slide: object,
        x: float,
        y: float,
        w: float,
        h: float,
        data: Sequence[Sequence[str]],
        header: bool = True,
        font_size: int = 10,
        emphasize_rows: Iterable[int] = (),
    ) -> None:
        rows, cols = len(data), len(data[0])
        shape = slide.shapes.add_table(rows, cols, Inches(x), Inches(y), Inches(w), Inches(h))
        tbl = shape.table
        for i, row in enumerate(data):
            for j, val in enumerate(row):
                cell = tbl.cell(i, j)
                cell.text = str(val)
                cell.margin_left = Inches(0.05)
                cell.margin_right = Inches(0.05)
                cell.margin_top = Inches(0.03)
                cell.margin_bottom = Inches(0.03)
                fill = cell.fill
                fill.solid()
                if header and i == 0:
                    fill.fore_color.rgb = COLORS["blue"]
                    font_color = RGBColor(255, 255, 255)
                    bold = True
                elif i in emphasize_rows:
                    fill.fore_color.rgb = RGBColor(231, 244, 244)
                    font_color = COLORS["ink"]
                    bold = True
                elif i % 2 == 0:
                    fill.fore_color.rgb = RGBColor(245, 247, 250)
                    font_color = COLORS["ink"]
                    bold = False
                else:
                    fill.fore_color.rgb = RGBColor(255, 255, 255)
                    font_color = COLORS["ink"]
                    bold = False
                for p in cell.text_frame.paragraphs:
                    p.alignment = PP_ALIGN.CENTER
                    for run in p.runs:
                        run.font.name = FONT
                        run.font.size = Pt(font_size)
                        run.font.bold = bold
                        run.font.color.rgb = font_color

    def image(self, slide: object, path: str | Path, x: float, y: float, w: float, h: float, label: str | None = None) -> bool:
        path = REPO_ROOT / path if not Path(path).is_absolute() else Path(path)
        if not path.exists():
            self.card(slide, x, y, w, h, "Missing figure", str(path.relative_to(REPO_ROOT)), COLORS["red"], 10, 11)
            return False
        with Image.open(path) as img:
            iw, ih = img.size
        scale = min(w / iw, h / ih)
        tw, th = iw * scale, ih * scale
        px, py = x + (w - tw) / 2, y + (h - th) / 2
        slide.shapes.add_picture(str(path), Inches(px), Inches(py), width=Inches(tw), height=Inches(th))
        if label:
            self.text(slide, x, y + h + 0.04, w, 0.22, label, 8, COLORS["muted"], align=PP_ALIGN.CENTER)
        return True

    def pill(self, slide: object, x: float, y: float, w: float, text: str, color: RGBColor) -> None:
        shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(0.36))
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()
        shape.adjustments[0] = 0.5
        tf = shape.text_frame
        tf.clear()
        tf.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = text
        r.font.name = FONT
        r.font.size = Pt(10)
        r.font.bold = True
        r.font.color.rgb = RGBColor(255, 255, 255)

    def save(self) -> None:
        OUT_PPTX.parent.mkdir(parents=True, exist_ok=True)
        self.prs.save(OUT_PPTX)
        lines = ["# RADIO-GS Submission Presentation Outline", ""]
        for i, (title, notes) in enumerate(self.outline, 1):
            lines.append(f"## {i}. {title}")
            if notes:
                lines.extend([f"- {n}" for n in notes])
            lines.append("")
        OUT_OUTLINE.write_text("\n".join(lines), encoding="utf-8")


def add_pipeline_diagram(deck: DeckBuilder, slide: object) -> None:
    stages = [
        ("RGB + Poses", "训练视角与冻结 3DGS 几何", COLORS["blue"]),
        ("RADIO Teacher", "C-RADIOv4-H dense features", COLORS["teal"]),
        ("Hybrid Gaussian Field", "per-Gaussian latent + coarse branch", COLORS["amber"]),
        ("HCD + Refiner", "compact codec + screen alignment", COLORS["violet"]),
        ("Rendered Features", "novel-view RADIO-compatible maps", COLORS["green"]),
    ]
    x0, y, bw, gap = 0.55, 1.65, 2.25, 0.23
    for idx, (name, body, color) in enumerate(stages):
        x = x0 + idx * (bw + gap)
        deck.card(slide, x, y, bw, 1.28, name, body, color, 10, 11)
        if idx < len(stages) - 1:
            deck.line(slide, x + bw + 0.02, y + 0.64, x + bw + gap - 0.03, y + 0.64, COLORS["muted"])
    deck.card(slide, 1.0, 4.0, 3.3, 1.15, "Training losses", "cos/L1 feature distill\nFDH depth head\nDINO/SAM adaptor consistency", COLORS["red"], 10, 12)
    deck.card(slide, 5.0, 4.0, 3.3, 1.15, "Inference", "render feature map once\nquery by text / adaptor task\nno task-specific finetuning", COLORS["green"], 10, 12)
    deck.card(slide, 9.0, 4.0, 3.0, 1.15, "Paper claim", "3D Gaussian scene as reusable foundation-feature memory", COLORS["blue"], 10, 12)


def add_metric_bar(deck: DeckBuilder, slide: object, x: float, y: float, label: str, value: float, max_value: float, color: RGBColor) -> None:
    deck.text(slide, x, y, 2.25, 0.25, label, 10, COLORS["ink"], bold=True)
    base = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x + 2.45), Inches(y + 0.03), Inches(3.3), Inches(0.16))
    base.fill.solid()
    base.fill.fore_color.rgb = COLORS["line"]
    base.line.fill.background()
    barw = 3.3 * min(max(value / max_value, 0), 1)
    bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x + 2.45), Inches(y + 0.03), Inches(barw), Inches(0.16))
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()
    deck.text(slide, x + 5.85, y - 0.02, 0.7, 0.25, f"{value:.3f}", 10, COLORS["ink"], bold=True, align=PP_ALIGN.RIGHT)


def build_deck() -> None:
    deck = DeckBuilder()

    slide = deck.add_slide("RADIO-GS", "Title")
    deck.text(slide, 0.55, 1.55, 11.4, 0.9, "Foundation Feature Reconstruction in 3D Gaussian Scenes", 30, COLORS["ink"], bold=True)
    deck.text(slide, 0.6, 2.55, 10.7, 0.45, "面向论文投稿的完整项目汇报", 17, COLORS["muted"])
    deck.card(slide, 0.75, 3.55, 3.1, 1.15, "Main LERF", "0.8712 LocAcc\n0.4941 mIoU", COLORS["blue"])
    deck.card(slide, 4.1, 3.55, 3.1, 1.15, "ScanNet v67", "0.3538 / 0.3573 / 0.4293 mIoU", COLORS["teal"], body_size=15)
    deck.card(slide, 7.45, 3.55, 3.1, 1.15, "Core Ablation", "HCD: -0.3272 LocAcc when removed", COLORS["amber"], body_size=15)
    deck.text(slide, 0.6, 6.35, 10.5, 0.3, "Current package: LaTeX draft + frozen reports + qualitative figures + component ablations", 11, COLORS["muted"])
    deck.add_note("开场用一句话定义：不是训练一个语言分类器，而是把 RADIO foundation features 重建成一个可渲染的 3D feature memory。")

    slide = deck.add_slide("一句话主张", "Story")
    deck.text(slide, 0.75, 1.45, 11.5, 0.55, "3D Gaussian scene 可以成为 reusable foundation-feature memory。", 26, COLORS["blue"], bold=True)
    deck.bullets(slide, 1.0, 2.35, 11.1, 2.4, [
        "输入：已有 3DGS 几何 + 多视角 RGB + 冻结 RADIO teacher features。",
        "训练：学习 hybrid Gaussian feature field，重建 novel-view RADIO-compatible feature maps。",
        "推理：不用重新跑 teacher，也不用训练任务分类器；渲染出的 feature 可以做 text grounding、ScanNet point query、SAM/DINO probes。",
        "投稿角度：方法创新、主结果、消融、定性和失败分析已经形成闭合逻辑。"
    ], 16)
    deck.add_note("这页给听众建立主线，后面所有方法和实验都围绕 feature memory 展开。")

    slide = deck.add_slide("动机：2D foundation features 与 3D deployment 的断层", "Motivation")
    deck.card(slide, 0.65, 1.45, 3.7, 1.55, "2D foundation models", "RADIO / DINO / SAM 特征强，但只在已有 RGB 帧上可直接获得", COLORS["blue"], 12, 15)
    deck.card(slide, 4.8, 1.45, 3.7, 1.55, "3D scene representation", "3DGS 擅长 novel-view RGB，但通常不保存可复用语义特征", COLORS["teal"], 12, 15)
    deck.card(slide, 8.95, 1.45, 3.7, 1.55, "Open-vocabulary query", "部署时需要从新视角回答文本、分割和匹配问题", COLORS["amber"], 12, 15)
    deck.text(slide, 0.9, 4.05, 11.2, 0.55, "核心问题：能否把冻结视觉 foundation model 的 dense features 重建进 3D Gaussian scene，并在 novel views 上保持下游可用性？", 20, COLORS["ink"], bold=True)
    deck.add_note("强调问题不是传统 3D 语义分割，也不是 2D feature extraction，而是 3D feature reconstruction。")

    slide = deck.add_slide("待解决的问题", "Problem")
    deck.bullets(slide, 0.8, 1.35, 5.8, 4.7, [
        "高维特征存储：RADIO spatial feature 维度高，直接贴到每个 Gaussian 代价高。",
        "多视角一致性：feature map 要在 novel view 下稳定，不只是训练视角拟合。",
        "任务兼容性：rendered feature 需要继续兼容 frozen SigLIP/DINO/SAM heads。",
        "定位峰稳定：LERF LocAcc 是 argmax-in-mask，mIoU 变好不一定保住 peak。",
        "论文可证性：需要闭合 protocol、provenance、消融和失败案例。"
    ], 15)
    deck.card(slide, 7.2, 1.55, 4.6, 3.1, "Design target", "不是 scene-specific classifier\n不是只优化 text heatmap\n而是保持一个统一的 RADIO-compatible 3D feature field", COLORS["blue"], 13, 19)
    deck.add_note("这页为后面的 HCD、hybrid、refiner、FDH 和 adaptor supervision 做铺垫。")

    slide = deck.add_slide("方法总览", "Method")
    add_pipeline_diagram(deck, slide)
    deck.add_note("这页可以作为论文 Figure 2 的口头解释版本。先讲训练流，再讲 inference 流。")

    slide = deck.add_slide("方法组件 1：Hybrid Gaussian Feature Field", "Method")
    deck.bullets(slide, 0.7, 1.35, 5.7, 4.7, [
        "fine path：per-Gaussian compact latent feature，经 splatting 得到 view feature。",
        "coarse path：从 3D 空间分支查询更平滑的 scene-level context。",
        "fusion：两路合成 compact feature map，兼顾局部细节与全局一致性。",
        "消融发现：w/o hybrid macro LocAcc 降至 0.839，但 mIoU 提高到 0.507，说明 explicit 分支会扩大区域响应但削弱部分场景 peak。"
    ], 15)
    deck.card(slide, 7.0, 1.55, 4.7, 1.3, "Key evidence", "Full 0.858 LocAcc / 0.485 mIoU\nw/o hybrid 0.839 / 0.507", COLORS["teal"], 12, 18)
    deck.card(slide, 7.0, 3.4, 4.7, 1.3, "Interpretation", "Hybrid 主要不是提升所有 mIoU，而是稳定定位峰和跨场景表现。", COLORS["amber"], 12, 16)
    deck.add_note("这页回答用户之前关心的 LocAcc 下降原因：区域覆盖和 peak localization 是两种不同指标。")

    slide = deck.add_slide("方法组件 2：HCD Codec 与 Screen-Space Refiner", "Method")
    deck.bullets(slide, 0.75, 1.35, 5.65, 4.7, [
        "HCD：非线性 bottleneck codec，将 compact latent 解码到 RADIO feature space。",
        "目的：避免直接 1x1 projection 对高维 feature manifold 的表达不足。",
        "Refiner：在 screen space 用 alpha/depth/RGB 等 guide 修正边界与局部对齐。",
        "消融：w/o HCD 是最大退化；w/o refiner 主要损失 LocAcc。"
    ], 15)
    rows = [
        ["Variant", "LocAcc", "mIoU", "Delta LocAcc"],
        ["Full", "0.8578", "0.4850", "0.0000"],
        ["w/o HCD", "0.5306", "0.2596", "-0.3272"],
        ["w/o refiner", "0.8401", "0.4796", "-0.0177"],
    ]
    deck.table(slide, 6.9, 1.55, 5.4, 1.75, rows, font_size=11, emphasize_rows=[1])
    deck.card(slide, 7.15, 4.05, 4.9, 1.2, "Takeaway", "HCD 是架构必要项；refiner 是较小但稳定的 peak/boundary 修正项。", COLORS["blue"], 12, 16)
    deck.add_note("强调 HCD 消融是最强证据，不是调参带来的小差别。")

    slide = deck.add_slide("方法组件 3：FDH Warm-Start 与冻结头监督", "Method")
    deck.bullets(slide, 0.7, 1.35, 5.6, 4.7, [
        "先训练 no-FDH feature field，再 warm-start 到 FDH refinement。",
        "冻结 depth head 只提供几何感知的 feature regularization，不变成任务专用预测器。",
        "避免早期 feature reconstruction 与几何监督冲突。",
        "消融：w/o FDH macro LocAcc 0.8018，full 为 0.8578。"
    ], 15)
    add_metric_bar(deck, slide, 7.0, 1.55, "Full", 0.8578, 0.9, COLORS["blue"])
    add_metric_bar(deck, slide, 7.0, 2.05, "w/o FDH", 0.8018, 0.9, COLORS["amber"])
    add_metric_bar(deck, slide, 7.0, 2.85, "Full mIoU", 0.4850, 0.55, COLORS["teal"])
    add_metric_bar(deck, slide, 7.0, 3.35, "w/o FDH mIoU", 0.4236, 0.55, COLORS["red"])
    deck.card(slide, 7.0, 4.45, 5.1, 1.0, "Conclusion", "+0.056 LocAcc / +0.061 mIoU：FDH 是强训练路线增益。", COLORS["green"], 12, 15)

    slide = deck.add_slide("方法组件 4：RADIO Adaptors", "Method")
    deck.bullets(slide, 0.7, 1.35, 5.65, 4.9, [
        "官方 RADIO adaptor：SigLIP2-g、DINOv3、SAM3。",
        "SigLIP2-g：主 text-grounding evaluator。",
        "DINOv3：relation / cross-view token affinity，增强结构和匹配一致性。",
        "SAM3：soft-region supervision，借鉴 SAM/FMongo-style region prior。",
        "原则：训练时作为 regularizer，推理时仍输出统一 RADIO-compatible rendered feature。"
    ], 15)
    deck.pill(slide, 7.0, 1.55, 1.6, "SigLIP2-g", COLORS["blue"])
    deck.pill(slide, 8.85, 1.55, 1.45, "DINOv3", COLORS["teal"])
    deck.pill(slide, 10.55, 1.55, 1.25, "SAM3", COLORS["amber"])
    deck.card(slide, 7.0, 2.35, 4.9, 1.2, "Promoted selector", "macro LocAcc 0.8712\nmacro mIoU 0.4979", COLORS["green"], 12, 18)
    deck.card(slide, 7.0, 4.0, 4.9, 1.35, "Caveat", "DINO/SAM 分支经常提升区域 overlap，但可能移动 argmax peak，因此只作为 ablation/diagnostic。", COLORS["red"], 12, 14)

    slide = deck.add_slide("实验协议与证据边界", "Experiments")
    rows = [
        ["Evidence", "Protocol", "Paper role"],
        ["LERF-OVS", "4 scenes, rendered-feature grounding", "Main table"],
        ["Rendered vs RADIO RGB", "same frozen SigLIP2 evaluator", "feature-memory claim"],
        ["Component ablation", "seed-7 controlled, 4 scenes", "core design proof"],
        ["ScanNet v67", "10-scene direct point query", "cross-domain support"],
        ["Efficiency/profile", "eval wall time + peak VRAM", "cost table"],
        ["External baselines", "identity verified, values unresolved", "must be caveated"],
    ]
    deck.table(slide, 0.65, 1.35, 12.0, 4.2, rows, font_size=10, emphasize_rows=[1, 3])
    deck.add_note("这页避免 reviewer 认为混用协议。强调 main table、ablation、diagnostic 的边界。")

    slide = deck.add_slide("主结果：LERF-OVS Rendered-Feature Grounding", "Results")
    rows = [
        ["Scene", "LocAcc", "mIoU", "Temp"],
        ["Figurines", "0.8214", "0.4308", "50"],
        ["Ramen", "0.9014", "0.5862", "40"],
        ["Teatime", "0.8983", "0.5486", "25"],
        ["Waldo Kitchen", "0.8636", "0.4106", "25"],
        ["Macro", "0.8712", "0.4941", "--"],
    ]
    deck.table(slide, 0.7, 1.35, 6.2, 2.6, rows, font_size=12, emphasize_rows=[5])
    deck.card(slide, 7.4, 1.55, 4.6, 1.1, "Main claim", "RADIO-GS 在四个 LERF-OVS 场景上形成强 open-vocabulary grounding 主结果。", COLORS["blue"], 12, 16)
    deck.card(slide, 7.4, 3.05, 4.6, 1.1, "Hardest case", "Figurines 仍是小目标/feature resolution failure analysis 的重点。", COLORS["amber"], 12, 16)
    deck.add_note("注意主表是 current-best freeze，不和 seed-7 ablation 混为一谈。")

    slide = deck.add_slide("Rendered Features vs. Original RADIO RGB", "Results")
    rows = [
        ["Source", "Macro LocAcc", "Macro mIoU"],
        ["Original RADIO RGB", "0.7985", "0.4922"],
        ["RADIO-GS rendered", "0.8712", "0.4941"],
        ["Gain", "+0.0727", "+0.0019"],
    ]
    deck.table(slide, 0.85, 1.4, 5.6, 1.95, rows, font_size=13, emphasize_rows=[2, 3])
    deck.text(slide, 7.0, 1.55, 5.2, 0.55, "关键推论", 22, COLORS["blue"], bold=True)
    deck.bullets(slide, 7.0, 2.25, 5.1, 2.7, [
        "Rendered feature 不是 teacher 的简单有损拷贝。",
        "多视角 3D reconstruction 可以 sharpen task-relevant response。",
        "这直接支撑“3D feature memory”的论文主张。"
    ], 16)
    deck.add_note("这页是最强 story 页之一：为什么重建的 feature 可能比单帧 teacher 更适合 novel-view grounding。")

    slide = deck.add_slide("定性结果：LERF grounding overlays", "Qualitative")
    deck.image(slide, "output/radio_gs/paper_figures/submission_freeze_lerf_qualitative_comparison.png", 0.55, 1.25, 12.2, 5.45)
    deck.add_note("用这页展示四个场景的 rendered feature heatmap 和视觉质量，强调来自 frozen shortlist。")

    slide = deck.add_slide("核心组件消融：闭合证据链", "Ablation")
    rows = [
        ["Variant", "LocAcc", "Delta", "mIoU", "Interpretation"],
        ["Full", "0.8578", "0.0000", "0.4850", "main architecture"],
        ["w/o FDH", "0.8018", "-0.0560", "0.4236", "geometry warm-start matters"],
        ["w/o refiner", "0.8401", "-0.0177", "0.4796", "peak/boundary correction"],
        ["w/o hybrid", "0.8394", "-0.0184", "0.5069", "region up, peak down"],
        ["w/o HCD", "0.5306", "-0.3272", "0.2596", "codec is essential"],
    ]
    deck.table(slide, 0.45, 1.25, 12.45, 2.65, rows, font_size=10, emphasize_rows=[1, 5])
    deck.bullets(slide, 0.8, 4.35, 11.5, 1.7, [
        "HCD 是最大必要项：移除后四场景全部明显下降。",
        "FDH 是最强训练路线增益：同时提升 LocAcc 与 mIoU。",
        "refiner 和 hybrid 主要解释 LocAcc/mIoU tradeoff：区域热图变好不等于 argmax 更准。"
    ], 15)

    slide = deck.add_slide("为什么 mIoU 提升但 LocAcc 下降？", "Failure")
    deck.card(slide, 0.8, 1.45, 3.55, 1.35, "mIoU", "看 thresholded region overlap\n允许热图更宽、更平滑", COLORS["teal"], 12, 16)
    deck.card(slide, 4.9, 1.45, 3.55, 1.35, "LocAcc", "看 single argmax 是否落入 mask\n一格偏移就可能失败", COLORS["amber"], 12, 16)
    deck.card(slide, 9.0, 1.45, 3.0, 1.35, "Figurines", "小目标 + 低 feature resolution\n最容易触发 peak shift", COLORS["red"], 12, 15)
    deck.bullets(slide, 1.0, 3.55, 10.8, 2.0, [
        "DINO/SAM 或 explicit 分支常让 object region 更连贯，因此 mIoU 上升。",
        "但 LERF LocAcc 只看最高点；peak 被平滑到边界外就会下降。",
        "refiner、hybrid、spatial text-heatmap 的价值在于保护 peak，而不只是扩大响应区域。"
    ], 16)

    slide = deck.add_slide("Adaptor ablations：DINOv3 / SAM3", "Ablation")
    rows = [
        ["Scene", "Promoted variant", "LocAcc", "mIoU"],
        ["Figurines", "DINO cv + spatial text", "0.8214", "0.4343"],
        ["Ramen", "DINO rel. + SAM3 region", "0.9014", "0.5873"],
        ["Teatime", "DINO rel. + SAM3 region", "0.8983", "0.5592"],
        ["Waldo", "baseline", "0.8636", "0.4106"],
        ["Macro", "selector", "0.8712", "0.4979"],
    ]
    deck.table(slide, 0.65, 1.25, 6.2, 2.7, rows, font_size=10, emphasize_rows=[5])
    deck.image(slide, "paper/figures/lerf_adaptor_downstream_qualitative.png", 7.2, 1.3, 5.2, 3.2)
    deck.bullets(slide, 0.9, 4.55, 11.2, 1.2, [
        "Promoted selector 只选择不损失 LocAcc 的 adaptor/cross-view checkpoints。",
        "目前是 ablation/diagnostic，不替换保守主表。"
    ], 15)

    slide = deck.add_slide("SAM/DINO 下游探针", "Diagnostics")
    rows = [
        ["Task", "Teacher", "Rendered", "Readout"],
        ["SAM3 point prompt", "1.000 / 0.370", "1.000 / 0.417", "mask mIoU win"],
        ["SAM3 box prompt", "0.870 / 0.656", "0.822 / 0.664", "mask mIoU win"],
        ["DINO dense matching", "0.590 / 0.854", "0.554 / 0.905", "smooth similarity"],
        ["DINO mask prop.", "0.716 / 0.392", "0.738 / 0.368", "gap narrowed"],
    ]
    deck.table(slide, 0.55, 1.25, 6.3, 2.35, rows, font_size=9, emphasize_rows=[3])
    deck.image(slide, "paper/figures/lerf_sam_dino_tasks_qualitative.png", 7.15, 1.25, 5.4, 2.75)
    deck.bullets(slide, 0.9, 4.4, 11.1, 1.45, [
        "Rendered features 在 SAM3-adaptor prompt mask mIoU 上超过 frame-wise teacher。",
        "DINO mask propagation 经 source-background contrast 明显改善，但 mIoU 仍未完全追平 teacher。"
    ], 15)

    slide = deck.add_slide("ScanNet v67：跨域 direct point-query", "Results")
    rows = [
        ["Split", "mIoU", "mAcc"],
        ["19 classes", "0.3538", "0.6076"],
        ["15 classes", "0.3573", "0.6203"],
        ["10 classes", "0.4293", "0.7051"],
    ]
    deck.table(slide, 0.75, 1.35, 4.8, 1.8, rows, font_size=13, emphasize_rows=[3])
    deck.image(slide, "output/radio_gs/paper_figures/vis_batch_20260501_overlay/contact_sheets/scannet_scene0070_00_gt_pred_error_split19.png", 6.0, 1.25, 6.6, 4.45)
    deck.bullets(slide, 0.9, 3.8, 4.6, 1.6, [
        "10-scene fair protocol。",
        "用于证明 reconstructed feature field 不是只对 LERF text grounding 有效。",
        "不包装成完整 ScanNet leaderboard。"
    ], 13)

    slide = deck.add_slide("ProFuse-inspired DINO cross-view diagnostics", "Diagnostics")
    rows = [
        ["Scene / Branch", "split19", "split15", "split10"],
        ["scene0070 baseline", "0.2297", "0.2405", "0.3238"],
        ["scene0070 DINO cv 0.001", "0.2437", "0.2466", "0.3284"],
        ["scene0645 baseline", "0.2381", "0.2458", "0.2875"],
        ["scene0645 DINO cv 0.003", "0.2427", "0.2500", "0.2833"],
    ]
    deck.table(slide, 0.65, 1.35, 7.2, 2.4, rows, font_size=11, emphasize_rows=[2, 4])
    deck.card(slide, 8.35, 1.55, 3.8, 1.4, "Interpretation", "DINO cross-view context 对部分 ScanNet 场景有效，但还不是完整 10-scene 主线。", COLORS["teal"], 12, 15)
    deck.card(slide, 8.35, 3.45, 3.8, 1.35, "Next step", "做 conservative-weight 全 10-scene sweep，再决定是否进入 main paper。", COLORS["amber"], 12, 15)

    slide = deck.add_slide("效率与成本：分开报告不同 measurement type", "Efficiency")
    rows = [
        ["Workload", "Wall time", "Peak VRAM"],
        ["LERF Figurines overlay", "26.198 s", "1568 MiB"],
        ["LERF Ramen overlay", "40.474 s", "1762 MiB"],
        ["LERF Teatime overlay", "36.997 s", "1850 MiB"],
        ["LERF Waldo overlay", "21.101 s", "2076 MiB"],
        ["ScanNet v67 10-scene eval", "150.903 s", "1666 MiB"],
    ]
    deck.table(slide, 0.65, 1.25, 6.8, 2.8, rows, font_size=11, emphasize_rows=[5])
    deck.bullets(slide, 8.0, 1.35, 4.4, 3.7, [
        "训练成本：来自 training.log，适合单独成表。",
        "eval latency：来自 profile workload。",
        "peak VRAM：来自 explicit GPU telemetry。",
        "投稿时避免把不同 measurement type 混在一个结论里。"
    ], 14)

    slide = deck.add_slide("Baseline provenance：投稿前必须保守", "Submission")
    rows = [
        ["Method", "Status", "Risk"],
        ["LERF", "identity verified", "repo values not official-matched"],
        ["LangSplat", "identity verified", "protocol mismatch across papers"],
        ["LEGaussians", "identity verified", "scene naming / supplement row mismatch"],
        ["RADIO-GS", "JSON-backed", "internal result frozen"],
    ]
    deck.table(slide, 0.65, 1.35, 7.4, 2.5, rows, font_size=11, emphasize_rows=[4])
    deck.card(slide, 8.55, 1.45, 3.7, 1.35, "P0 decision", "不混用来源不闭合的外部数值；官方锚定或统一协议复现。", COLORS["red"], 12, 15)
    deck.card(slide, 8.55, 3.35, 3.7, 1.35, "Current use", "外部 baseline 只作为 draft comparison caveat，不作为 freeze-safe final claim。", COLORS["amber"], 12, 15)

    slide = deck.add_slide("项目完成度与投稿主线", "Submission")
    deck.card(slide, 0.75, 1.35, 3.2, 1.25, "Completion", "约 90%\n投稿包已成型", COLORS["green"])
    deck.card(slide, 4.25, 1.35, 3.2, 1.25, "Mainline", "LERF-OVS + ScanNet v67\nfeature memory story", COLORS["blue"], body_size=17)
    deck.card(slide, 7.75, 1.35, 3.8, 1.25, "Risks", "external baseline / final cost table / venue polish", COLORS["amber"], body_size=16)
    deck.bullets(slide, 1.0, 3.35, 10.8, 2.5, [
        "主线只保留：full RADIO-GS、same-evaluator teacher-vs-rendered、核心组件消融、ScanNet v67、qualitative + profile。",
        "DINO/SAM/ProFuse branches 放入 ablation/diagnostic，不抢主贡献。",
        "历史分支归档，仅用于解释设计选择和 failure analysis。"
    ], 15)

    slide = deck.add_slide("论文结构建议", "Paper")
    rows = [
        ["Section", "Main message"],
        ["Introduction", "3DGS as reusable foundation-feature memory"],
        ["Related Work", "language fields / 3DGS / RADIO / SAM-DINO regularization"],
        ["Method", "hybrid field + HCD + refiner + FDH + adaptor consistency"],
        ["Experiments", "LERF main, teacher-vs-rendered, ablations, ScanNet, efficiency"],
        ["Failure Analysis", "Figurines, feature resolution, peak-vs-region tradeoff"],
        ["Limitations", "baseline provenance, SAM gap, full ScanNet leaderboard not claimed"],
    ]
    deck.table(slide, 0.75, 1.25, 11.8, 4.35, rows, font_size=11, emphasize_rows=[1, 4])

    slide = deck.add_slide("投稿前剩余工作", "Next")
    deck.bullets(slide, 0.85, 1.35, 11.5, 4.85, [
        "P0 baseline：统一协议复现或只保留官方可锚定数值并写清 protocol mismatch。",
        "P2 efficiency：把训练成本、eval latency、peak VRAM 分三类表格写入论文。",
        "P3 failure：围绕 Figurines 写小目标尺寸、feature resolution、peak shift 和 refiner 效果。",
        "图表：把 framework figure、qualitative comparison、ablation table 做到最终版式。",
        "文字：迁移到目标会议模板，压缩相关工作，统一 terminology 和 artifact provenance。"
    ], 16)

    slide = deck.add_slide("Takeaways", "Close")
    deck.text(slide, 0.8, 1.4, 11.5, 0.5, "RADIO-GS 的投稿主张已经自洽：", 24, COLORS["blue"], bold=True)
    deck.bullets(slide, 1.0, 2.25, 10.8, 3.1, [
        "方法上：将 RADIO foundation features 重建进 3D Gaussian scene，形成可渲染 feature memory。",
        "实验上：LERF main result、teacher-vs-rendered、ScanNet transfer 和 component ablations 已经闭合。",
        "消融上：HCD / FDH / refiner / hybrid 的作用可以用同一 LERF seed-7 protocol 解释。",
        "投稿上：剩余风险主要是外部 baseline provenance 和最终写作/版式，而不是核心方法完成度。"
    ], 17)
    deck.add_note("最后把问题收回投稿目标：方法和证据链已经足够，下一步是 freeze-safe presentation and paper polish。")

    slide = deck.add_slide("Backup：详细 LERF component table", "Backup")
    rows = [
        ["Variant", "Fig.", "Ramen", "Tea.", "Waldo", "LocAcc", "mIoU"],
        ["Full", "0.768", "0.901", "0.898", "0.864", "0.858", "0.485"],
        ["w/o FDH", "0.696", "0.845", "0.847", "0.818", "0.802", "0.424"],
        ["w/o refiner", "0.768", "0.859", "0.915", "0.818", "0.840", "0.480"],
        ["w/o hybrid", "0.768", "0.873", "0.898", "0.818", "0.839", "0.507"],
        ["w/o HCD", "0.268", "0.648", "0.661", "0.545", "0.531", "0.260"],
    ]
    deck.table(slide, 0.55, 1.35, 12.2, 2.7, rows, font_size=11, emphasize_rows=[1, 5])
    deck.bullets(slide, 0.8, 4.55, 11.4, 1.2, [
        "Protocol: controlled seed-7 LERF-OVS rendered-feature temperature sweep。",
        "完整 provenance: output/radio_gs/reports/lerf_component_ablation.md / .json。"
    ], 14)

    slide = deck.add_slide("Backup：artifact map", "Backup")
    rows = [
        ["Artifact", "Path"],
        ["LaTeX draft", "paper/radio_gs_draft.tex"],
        ["Component ablation", "output/radio_gs/reports/lerf_component_ablation.md"],
        ["Freeze report", "output/radio_gs/reports/submission_freeze_report.md"],
        ["Qualitative shortlist", "output/radio_gs/reports/submission_freeze_figure_shortlist.md"],
        ["Efficiency profile", "output/radio_gs/reports/efficiency_profile.md"],
        ["Baseline audit", "output/radio_gs/reports/baseline_source_verification.md"],
    ]
    deck.table(slide, 0.65, 1.25, 12.0, 4.2, rows, font_size=10, emphasize_rows=[1, 2])

    deck.save()


if __name__ == "__main__":
    build_deck()
    print(f"Wrote {OUT_PPTX}")
    print(f"Wrote {OUT_OUTLINE}")
