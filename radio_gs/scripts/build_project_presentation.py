#!/usr/bin/env python3
"""Build a Chinese project presentation deck for GaussFM / RADIO-GS.

The deck is intentionally generated from a small script so the paper-facing
numbers, figure paths, and terminology stay reproducible.
"""

from __future__ import annotations

from pathlib import Path
from typing import Iterable, Sequence

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_AUTO_SIZE, MSO_VERTICAL_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


REPO_ROOT = Path(__file__).resolve().parents[2]
PAPER_DIR = REPO_ROOT / "paper"
FIG_DIR = PAPER_DIR / "figures"
OUT_PPTX = PAPER_DIR / "radio_gs_project_presentation_20260612.pptx"
OUT_OUTLINE = PAPER_DIR / "radio_gs_project_presentation_20260612_outline.md"

SLIDE_W = 13.333
SLIDE_H = 7.5

FONT = "Microsoft YaHei"

COLORS = {
    "ink": RGBColor(24, 31, 42),
    "muted": RGBColor(83, 96, 111),
    "paper": RGBColor(248, 250, 252),
    "panel": RGBColor(255, 255, 255),
    "line": RGBColor(213, 220, 228),
    "teal": RGBColor(17, 128, 127),
    "blue": RGBColor(44, 93, 165),
    "amber": RGBColor(204, 128, 34),
    "red": RGBColor(184, 68, 68),
    "green": RGBColor(75, 143, 94),
    "violet": RGBColor(95, 88, 170),
    "dark": RGBColor(18, 44, 57),
    "soft_teal": RGBColor(230, 246, 245),
    "soft_blue": RGBColor(232, 239, 252),
    "soft_amber": RGBColor(252, 244, 230),
}


class Deck:
    def __init__(self) -> None:
        self.prs = Presentation()
        self.prs.slide_width = Inches(SLIDE_W)
        self.prs.slide_height = Inches(SLIDE_H)
        self.blank = self.prs.slide_layouts[6]
        self.outline: list[tuple[str, list[str]]] = []

    def save(self) -> None:
        OUT_PPTX.parent.mkdir(parents=True, exist_ok=True)
        self.prs.save(OUT_PPTX)
        lines = [
            "# GaussFM / RADIO-GS 项目汇报 PPT 大纲",
            "",
            f"- 输出文件：`{OUT_PPTX}`",
            "- 口径：导师/项目中期汇报 + 顶刊投稿准备",
            "",
        ]
        for i, (title, notes) in enumerate(self.outline, 1):
            lines.append(f"## {i}. {title}")
            for note in notes:
                lines.append(f"- {note}")
            lines.append("")
        OUT_OUTLINE.write_text("\n".join(lines), encoding="utf-8")

    def slide(self, title: str, section: str | None = None) -> object:
        s = self.prs.slides.add_slide(self.blank)
        s.background.fill.solid()
        s.background.fill.fore_color.rgb = COLORS["paper"]
        if section:
            self.text(s, 0.50, 0.18, 4.0, 0.24, section.upper(), 8, COLORS["teal"], bold=True)
        self.text(s, 0.50, 0.43, 11.2, 0.45, title, 22, COLORS["ink"], bold=True)
        self.line(s, 0.50, 1.00, 12.75, 1.00, COLORS["line"])
        self.footer(s)
        self.outline.append((title, []))
        return s

    def note(self, text: str) -> None:
        self.outline[-1][1].append(text)

    def footer(self, s: object) -> None:
        idx = len(self.prs.slides)
        self.text(s, 0.50, 7.08, 4.0, 0.18, "GaussFM / RADIO-GS Project Presentation", 7, COLORS["muted"])
        self.text(s, 12.15, 7.08, 0.55, 0.18, str(idx), 7, COLORS["muted"], align=PP_ALIGN.RIGHT)

    def text(
        self,
        s: object,
        x: float,
        y: float,
        w: float,
        h: float,
        value: str,
        size: int,
        color: RGBColor | None = None,
        bold: bool = False,
        align: PP_ALIGN | None = None,
    ) -> object:
        box = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = box.text_frame
        tf.clear()
        tf.word_wrap = True
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        tf.margin_left = Inches(0.02)
        tf.margin_right = Inches(0.02)
        tf.margin_top = Inches(0.01)
        tf.margin_bottom = Inches(0.01)
        p = tf.paragraphs[0]
        if align is not None:
            p.alignment = align
        r = p.add_run()
        r.text = value
        r.font.name = FONT
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color or COLORS["ink"]
        return box

    def bullets(
        self,
        s: object,
        x: float,
        y: float,
        w: float,
        h: float,
        items: Sequence[str],
        size: int = 14,
        color: RGBColor | None = None,
    ) -> None:
        box = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = box.text_frame
        tf.clear()
        tf.word_wrap = True
        tf.margin_left = Inches(0.12)
        tf.margin_right = Inches(0.04)
        tf.margin_top = Inches(0.02)
        for i, item in enumerate(items):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = item
            p.level = 0
            p.space_after = Pt(8)
            p.font.name = FONT
            p.font.size = Pt(size)
            p.font.color.rgb = color or COLORS["ink"]

    def card(
        self,
        s: object,
        x: float,
        y: float,
        w: float,
        h: float,
        title: str,
        body: str,
        accent: RGBColor,
        body_size: int = 16,
    ) -> None:
        shape = s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
        shape.fill.solid()
        shape.fill.fore_color.rgb = COLORS["panel"]
        shape.line.color.rgb = COLORS["line"]
        shape.adjustments[0] = 0.06
        bar = s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x), Inches(y), Inches(0.08), Inches(h))
        bar.fill.solid()
        bar.fill.fore_color.rgb = accent
        bar.line.fill.background()
        self.text(s, x + 0.20, y + 0.15, w - 0.35, 0.28, title, 11, COLORS["muted"], bold=True)
        self.text(s, x + 0.20, y + 0.52, w - 0.35, h - 0.60, body, body_size, COLORS["ink"], bold=True)

    def pill(self, s: object, x: float, y: float, w: float, text: str, color: RGBColor) -> None:
        shape = s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(0.34))
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()
        shape.adjustments[0] = 0.5
        tf = shape.text_frame
        tf.clear()
        tf.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = text
        r.font.name = FONT
        r.font.size = Pt(9)
        r.font.bold = True
        r.font.color.rgb = RGBColor(255, 255, 255)

    def line(self, s: object, x1: float, y1: float, x2: float, y2: float, color: RGBColor) -> None:
        c = s.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
        c.line.color.rgb = color
        c.line.width = Pt(1.3)

    def table(
        self,
        s: object,
        x: float,
        y: float,
        w: float,
        h: float,
        data: Sequence[Sequence[str]],
        font_size: int = 10,
        emphasize_rows: Iterable[int] = (),
    ) -> None:
        rows, cols = len(data), len(data[0])
        shape = s.shapes.add_table(rows, cols, Inches(x), Inches(y), Inches(w), Inches(h))
        tbl = shape.table
        for i, row in enumerate(data):
            for j, val in enumerate(row):
                cell = tbl.cell(i, j)
                cell.text = str(val)
                cell.margin_left = Inches(0.04)
                cell.margin_right = Inches(0.04)
                cell.margin_top = Inches(0.02)
                cell.margin_bottom = Inches(0.02)
                cell.fill.solid()
                if i == 0:
                    cell.fill.fore_color.rgb = COLORS["dark"]
                    fcolor = RGBColor(255, 255, 255)
                    bold = True
                elif i in emphasize_rows:
                    cell.fill.fore_color.rgb = COLORS["soft_teal"]
                    fcolor = COLORS["ink"]
                    bold = True
                elif i % 2 == 0:
                    cell.fill.fore_color.rgb = RGBColor(244, 247, 250)
                    fcolor = COLORS["ink"]
                    bold = False
                else:
                    cell.fill.fore_color.rgb = RGBColor(255, 255, 255)
                    fcolor = COLORS["ink"]
                    bold = False
                for p in cell.text_frame.paragraphs:
                    p.alignment = PP_ALIGN.CENTER
                    for r in p.runs:
                        r.font.name = FONT
                        r.font.size = Pt(font_size)
                        r.font.bold = bold
                        r.font.color.rgb = fcolor

    def image(self, s: object, path: str | Path, x: float, y: float, w: float, h: float, label: str | None = None) -> bool:
        p = Path(path)
        if not p.is_absolute():
            p = REPO_ROOT / p
        if not p.exists():
            self.card(s, x, y, w, h, "Missing figure", str(p), COLORS["red"], 10)
            return False
        with Image.open(p) as im:
            iw, ih = im.size
        scale = min(w / iw, h / ih)
        tw, th = iw * scale, ih * scale
        px, py = x + (w - tw) / 2, y + (h - th) / 2
        s.shapes.add_picture(str(p), Inches(px), Inches(py), width=Inches(tw), height=Inches(th))
        if label:
            self.text(s, x, y + h + 0.04, w, 0.20, label, 8, COLORS["muted"], align=PP_ALIGN.CENTER)
        return True

    def flow_node(
        self,
        s: object,
        x: float,
        y: float,
        w: float,
        h: float,
        title: str,
        body: str,
        fill: RGBColor,
        border: RGBColor,
    ) -> None:
        shape = s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
        shape.line.color.rgb = border
        shape.adjustments[0] = 0.07
        self.text(s, x + 0.15, y + 0.14, w - 0.30, 0.25, title, 11, border, bold=True)
        self.text(s, x + 0.15, y + 0.47, w - 0.30, h - 0.55, body, 12, COLORS["ink"])


def build() -> None:
    d = Deck()

    # 1
    s = d.prs.slides.add_slide(d.blank)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = COLORS["paper"]
    d.text(s, 0.62, 0.62, 11.8, 0.42, "GaussFM：紧凑基础特征 Gaussian 记忆", 26, COLORS["ink"], bold=True)
    d.text(
        s,
        0.62,
        1.15,
        11.2,
        0.54,
        "面向开放词汇 3D 场景理解的 compact RADIO feature memory",
        20,
        COLORS["teal"],
        bold=True,
    )
    d.text(s, 0.62, 1.95, 10.8, 0.60, "项目汇报 / 顶刊投稿准备\n2026-06-12", 14, COLORS["muted"])
    d.card(s, 0.68, 3.12, 3.75, 1.35, "核心目标", "把 frame-wise RADIO 特征压入可查询的 3D Gaussian 记忆", COLORS["teal"], 16)
    d.card(s, 4.82, 3.12, 3.75, 1.35, "核心优势", "低维存储、跨视角融合、二维/三维开放词汇任务统一支撑", COLORS["blue"], 16)
    d.card(s, 8.96, 3.12, 3.75, 1.35, "核心结果", "LERF 2D、LERF Direct3D、ScanNet VALA8 三条主线均有定量闭环", COLORS["amber"], 16)
    d.text(s, 0.70, 6.68, 11.5, 0.22, "注意：DINO/SigLIP2/SAM 在本文中主要作为 RADIO adaptor-space consistency 或 frozen-head evaluation，不是独立待重建的多教师特征。", 10, COLORS["muted"])
    d.outline.append(("标题与一句话主张", [
        "开场先定义本文不是做一个新的 3DGS RGB 重建，而是在冻结几何上学习 compact RADIO feature memory。",
        "用一句话讲清楚：小存储、可重建、可查询，并支撑 2D/3D 开放词汇理解。",
    ]))

    # 2
    s = d.slide("汇报主线：从问题到投稿闭环", "Roadmap")
    d.card(s, 0.70, 1.38, 2.25, 1.10, "1. 背景缺口", "2D foundation 特征强，但不能直接成为可部署 3D 场景记忆", COLORS["blue"], 13)
    d.card(s, 3.15, 1.38, 2.25, 1.10, "2. 方法设计", "冻结 RGB 3DGS 几何，学习 compact RADIO latent + hybrid context", COLORS["teal"], 13)
    d.card(s, 5.60, 1.38, 2.25, 1.10, "3. 查询机制", "按需重建 RADIO-compatible 特征，支持 2D/3D/point query", COLORS["green"], 13)
    d.card(s, 8.05, 1.38, 2.25, 1.10, "4. 实验闭环", "三 benchmark + 原始 RADIO 对比 + 存储/效率 + 消融", COLORS["amber"], 13)
    d.card(s, 10.50, 1.38, 2.25, 1.10, "5. 投稿计划", "TPAMI 叙事、图表布局、风险边界和下一步补强", COLORS["violet"], 13)
    d.bullets(
        s,
        1.05,
        3.08,
        11.4,
        2.7,
        [
            "讲故事顺序：不是“加了很多模块”，而是“把高维 frame-wise 特征压成可查询 3D 记忆”。",
            "方法命名建议：Compact Foundation-Feature Gaussian Memory；避免过度强调 readout 这个词。",
            "实验组织：主对比结果先行，消融和诊断支撑贡献点，qualitative 展示 2D 与 3D 查询差异。",
        ],
        16,
    )
    d.note("这一页用于给导师建立总路线，后面所有章节都围绕这个闭环展开。")

    # 3
    s = d.slide("背景：2D foundation features 与 3D 场景记忆之间有断层", "Motivation")
    d.flow_node(s, 0.80, 1.38, 3.20, 1.42, "Frame-wise RADIO", "每张图可得到强语义特征\n但视角间不一致、不能直接 3D 查询", COLORS["soft_blue"], COLORS["blue"])
    d.flow_node(s, 5.05, 1.38, 3.20, 1.42, "Explicit 3D feature storage", "直接存 1280-D per-Gaussian 特征\n存储大、噪声和碎片化明显", COLORS["soft_amber"], COLORS["amber"])
    d.flow_node(s, 9.30, 1.38, 3.20, 1.42, "Open-vocabulary 3D tasks", "既要 rendered-view 2D OVS\n也要 Gaussian/point-level direct query", COLORS["soft_teal"], COLORS["teal"])
    d.line(s, 4.15, 2.08, 4.85, 2.08, COLORS["line"])
    d.line(s, 8.40, 2.08, 9.10, 2.08, COLORS["line"])
    d.bullets(
        s,
        0.95,
        3.35,
        11.6,
        2.3,
        [
            "现有 2D VLM / RADIO 特征表达强，但它们是逐帧、逐像素结果；直接用于三维场景会带来跨视角漂移和存储冗余。",
            "现有 3DGS 开放词汇方法多在 2D heatmap 或 per-primitive language feature 上做文章，primitive-level 稳定性和小物体支持仍是难点。",
            "本文问题：能否学习一个紧凑、可重建、可按需查询的 3D feature memory，同时改善 2D 和 3D 下游任务？",
        ],
        15,
    )
    d.note("这里避免说别人完全没有 2D/3D 双接口，重点放在 compact + reconstructed feature quality + support stability。")

    # 4
    s = d.slide("核心假设：压缩不是损失信息，而是多视角去噪与重组", "Key Hypothesis")
    d.card(s, 0.90, 1.40, 3.65, 1.40, "高维特征压缩", "1280-D RADIO\n→ 32-D per-Gaussian latent\n→ 192-D compact scene code", COLORS["blue"], 17)
    d.card(s, 4.85, 1.40, 3.65, 1.40, "多视角证据注册", "把多视角一致语义压回 primitive field，而非推理时查 cache", COLORS["teal"], 17)
    d.card(s, 8.80, 1.40, 3.65, 1.40, "按需重建", "查询时由 compact memory + decoder 重建任务所需 feature", COLORS["amber"], 17)
    d.text(s, 1.00, 3.36, 11.3, 0.48, "因此，本文的贡献不是“简单低维化”，而是把 frame-wise foundation feature 转化为一个可部署的三维特征记忆。", 18, COLORS["ink"], bold=True)
    d.bullets(
        s,
        1.18,
        4.25,
        10.9,
        1.55,
        [
            "存储角度：per-Gaussian latent 是主要随场景规模增长的语义负载，显著小于显式 1280-D feature。",
            "性能角度：多视角融合后的 reconstructed scene features 在多个 frozen-head 任务上超过原始 frame-wise RADIO。",
        ],
        15,
    )
    d.note("这一页用于强调 compact 不是退化，是多视角信息瓶颈，有新意。")

    # 5
    s = d.slide("整体流程：先 RGB 3DGS 几何，再学习 compact RADIO feature memory", "Framework")
    d.flow_node(s, 0.72, 1.33, 2.55, 1.15, "Stage 0", "RGB images + cameras\n训练/加载 RGB 3DGS", COLORS["soft_blue"], COLORS["blue"])
    d.flow_node(s, 3.70, 1.33, 2.55, 1.15, "Frozen scaffold", "冻结 Gaussian 几何\nxyz / opacity / scale / rotation", COLORS["soft_teal"], COLORS["teal"])
    d.flow_node(s, 6.68, 1.33, 2.55, 1.15, "Feature learning", "学习 compact latent\n重建 RADIO-compatible features", COLORS["soft_amber"], COLORS["amber"])
    d.flow_node(s, 9.66, 1.33, 2.55, 1.15, "Task queries", "开放词汇 2D / 3D / ScanNet\n按需查询", RGBColor(238, 236, 252), COLORS["violet"])
    for x1, x2 in [(3.28, 3.62), (6.26, 6.60), (9.24, 9.58)]:
        d.line(s, x1, 1.91, x2, 1.91, COLORS["line"])
    d.bullets(
        s,
        0.92,
        3.08,
        11.8,
        2.20,
        [
            "本文实际主线：不是从零做 RGB supervision / RGB reconstruction，而是在已有 RGB 3DGS 几何上学习语义特征场。",
            "RGB 在 feature stage 中可作为 screen-space refiner guide 或边界辅助，但不是主要 feature target。",
            "训练监督主干来自 frozen RADIO features；DINO/SAM/SigLIP2 是 adaptor-space consistency 或下游评估接口。",
        ],
        15,
    )
    d.note("这是纠正框架图误解的关键页。")

    # 6
    s = d.slide("方法一：Hybrid Compact RADIO Feature Field", "Method")
    d.flow_node(s, 0.65, 1.28, 2.35, 1.22, "Fine branch", "32-D per-Gaussian latent\n通过 splatting 形成视图 compact map", COLORS["soft_blue"], COLORS["blue"])
    d.flow_node(s, 3.35, 1.28, 2.35, 1.22, "Coarse branch", "depth unprojection\nmulti-res spatial hash context", COLORS["soft_teal"], COLORS["teal"])
    d.flow_node(s, 6.05, 1.28, 2.35, 1.22, "Decoupled fusion", "semantic / geometry gates\nconfidence-aware fusion", COLORS["soft_amber"], COLORS["amber"])
    d.flow_node(s, 8.75, 1.28, 2.35, 1.22, "Compact code", "192-D scene feature code\n供 HCD decoder 重建", RGBColor(238, 236, 252), COLORS["violet"])
    d.flow_node(s, 10.95, 3.05, 1.45, 1.22, "RADIO", "1280-D\ncompatible", COLORS["panel"], COLORS["red"])
    for x1, x2 in [(3.05, 3.30), (5.75, 6.00), (8.45, 8.70), (10.50, 10.90)]:
        d.line(s, x1, 1.88, x2, 1.88, COLORS["line"])
    d.bullets(
        s,
        0.88,
        3.35,
        9.45,
        2.15,
        [
            "Fine path 保留 per-Gaussian 局部语义；coarse path 提供空间上下文，缓解 isolated primitive fragmentation。",
            "Decoupled fusion 让几何边界、可见性和语义置信度进入 compact feature，而不是只做 loss weighting。",
            "HCD codec 将 compact code 解码回 RADIO-compatible 1280-D feature，兼容 text/adaptor/frozen-head evaluation。",
        ],
        14,
    )
    d.note("这里回答用户关心的：compact memory learning 不只是一个 decoder，而是 fine latent + spatial hash + fusion + HCD codec。")

    # 7
    s = d.slide("方法二：Multiview Primitive Registration 压回 compact field", "Method")
    d.flow_node(s, 0.78, 1.38, 2.50, 1.18, "Multi-view evidence", "训练视角 RADIO feature\n可见性 / 置信度 / 边界", COLORS["soft_blue"], COLORS["blue"])
    d.flow_node(s, 3.70, 1.38, 2.50, 1.18, "Primitive registration", "把跨视角证据注册到 Gaussian primitives", COLORS["soft_teal"], COLORS["teal"])
    d.flow_node(s, 6.62, 1.38, 2.50, 1.18, "Compact distillation", "注册证据不作为推理 cache\n而是变成训练监督", COLORS["soft_amber"], COLORS["amber"])
    d.flow_node(s, 9.54, 1.38, 2.50, 1.18, "Deployable field", "推理只读 compact field\n无 VPR cache 依赖", RGBColor(238, 236, 252), COLORS["violet"])
    d.bullets(
        s,
        0.98,
        3.10,
        11.2,
        2.25,
        [
            "MPR 的角色：将多视角注册证据压缩为 primitive-level feature quality，而不是在推理时保留一个外部 cache。",
            "Direct 3D 的主要难点不是 feature cosine，而是 score-to-mask support；因此加入 support-calibrated primitive selection。",
            "Waldo 小物体 / 碎片化 / 低可见区域通过 prompt ensemble + score-component guard 得到明显恢复。",
        ],
        15,
    )
    d.note("这页把 VPR cache 风险改写成训练机制，而不是方法依赖。")

    # 8
    s = d.slide("训练目标：主干 RADIO 重建 + 任务相关一致性", "Method")
    d.table(
        s,
        0.80,
        1.35,
        11.75,
        3.25,
        [
            ["目标", "作用", "是否主干"],
            ["RADIO feature reconstruction", "重建 frozen RADIO 1280-D feature，形成 RADIO-compatible scene memory", "是"],
            ["Compact HCD target", "让低维 code 保留可解码的高维语义结构", "是"],
            ["Visibility / boundary / depth guides", "修正 splatting 与边界混合带来的 feature artifacts", "辅助"],
            ["SigLIP2 / SAM / DINO adaptor-space consistency", "验证/约束 RADIO-derived feature 对 frozen-head tasks 的可用性", "辅助"],
        ],
        font_size=10,
        emphasize_rows=[1, 2],
    )
    d.text(
        s,
        0.92,
        5.10,
        11.4,
        0.55,
        "表述边界：本文不声称同时重建真实 DINO/SAM/SigLIP2 encoder features；它们是 RADIO feature 经 adaptor 后的额外监督和下游可用性验证。",
        15,
        COLORS["red"],
        bold=True,
    )
    d.note("这是答辩/汇报中最容易被问到的概念边界。")

    # 9
    s = d.slide("推理/评估接口：同一 compact memory 支撑三类开放词汇任务", "Method")
    d.flow_node(s, 0.72, 1.32, 2.75, 1.35, "Rendered-view OVS", "render compact feature map\n→ reconstruct RADIO feature\n→ text heatmap / mask", COLORS["soft_blue"], COLORS["blue"])
    d.flow_node(s, 4.02, 1.32, 2.75, 1.35, "Direct 3D OVS", "query primitive scores\n→ support-calibrated selection\n→ render selected Gaussians", COLORS["soft_teal"], COLORS["teal"])
    d.flow_node(s, 7.32, 1.32, 2.75, 1.35, "ScanNet point query", "point/vertex feature query\n→ contextual kNN / text class query", COLORS["soft_amber"], COLORS["amber"])
    d.flow_node(s, 10.62, 1.32, 1.75, 1.35, "2D probes", "SAM/DINO/SigLIP\nfrozen-head tests", RGBColor(238, 236, 252), COLORS["violet"])
    d.bullets(
        s,
        0.92,
        3.22,
        11.5,
        2.2,
        [
            "建议 PPT 里少用 readout，改说“按需查询机制”或“任务接口”，并强调它们共享同一个 compact feature memory。",
            "二维和三维流程不同：2D 是先渲染 dense feature map 再 query；3D 是先在 primitive 上选择，再渲染 selected primitives 做评估。",
            "这能自然解释为什么定性图中 2D OVS 应是 RGB 上 heatmap/mask，3D OVS 应是白底 selected Gaussian / mask render。",
        ],
        15,
    )
    d.note("这页也是后续 qualitative 排版的逻辑依据。")

    # 10
    s = d.slide("实验设计：四条定量证据链", "Experiments")
    d.table(
        s,
        0.70,
        1.30,
        11.95,
        4.15,
        [
            ["证据链", "数据 / 协议", "证明什么", "主指标"],
            ["LERF rendered-view OVS", "4 scenes, 2D rendered feature maps", "重建 feature map 的开放词汇定位能力", "LocAcc / mIoU"],
            ["LERF direct 3D OVS", "OpenGaussian-style query-select-render", "primitive-level 直接查询与 object support", "mIoU / Acc@0.25"],
            ["ScanNet VALA8 point query", "VALA-aligned 8 scenes", "跨数据集点云开放词汇理解", "mIoU / mAcc"],
            ["Feature usability probes", "frozen SAM/DINO/SigLIP2 task heads", "重建特征是否优于 frame-wise RADIO", "mIoU / score"],
        ],
        font_size=10,
        emphasize_rows=[1, 2, 3, 4],
    )
    d.text(s, 0.92, 5.90, 11.0, 0.28, "消融实验另行回答：哪些模块真正贡献最大；storage/latency 回答部署价值。", 13, COLORS["muted"])
    d.note("四条证据链是论文实验章节的骨架。")

    # 11
    s = d.slide("主结果一：LERF rendered-view open-vocabulary grounding", "Results")
    d.table(
        s,
        0.75,
        1.32,
        7.25,
        2.40,
        [
            ["Method", "Fig.", "Ramen", "Tea.", "Waldo", "Macro"],
            ["LERF", "0.795", "0.625", "0.938", "0.815", "0.793"],
            ["LangSplat", "0.804", "0.732", "0.881", "0.955", "0.843"],
            ["LEGaussians", "0.767", "0.737", "0.683", "0.523", "0.678"],
            ["GaussFM", "0.821", "0.901", "0.898", "0.818", "0.860"],
        ],
        font_size=10,
        emphasize_rows=[4],
    )
    d.card(s, 8.35, 1.34, 3.85, 1.05, "补充控制", "Frame-wise RADIO mIoU 0.4634\nGaussFM rendered mIoU 0.5889", COLORS["teal"], 17)
    d.card(s, 8.35, 2.72, 3.85, 1.05, "结论", "多视角 compact memory 不只是复现 RADIO，而是提升 rendered-view 下游可用性", COLORS["blue"], 15)
    d.image(s, FIG_DIR / "lerf_rendered_grounding_qualitative.png", 0.90, 4.20, 11.55, 1.95, "Rendered-view grounding qualitative")
    d.note("这里先放 LERF rendered 主表，说明 2D 查询能力。")

    # 12
    s = d.slide("主结果二：LERF direct 3D object selection", "Results")
    d.table(
        s,
        0.70,
        1.30,
        8.05,
        3.35,
        [
            ["Method", "mIoU", "Acc@0.25", "备注"],
            ["OpenGaussian", "0.384", "0.514", "official context"],
            ["Dr. Splat", "0.433", "0.643", "published context"],
            ["GaussFM MPR", "0.480", "0.676", "registered evidence diagnostic"],
            ["GaussFM compact", "0.501", "0.704", "no VPR cache / no official SAM readout"],
            ["GaussFM + SAM3 box", "0.570", "0.684", "assisted boundary readout"],
        ],
        font_size=10,
        emphasize_rows=[4],
    )
    d.card(s, 9.15, 1.38, 3.10, 1.35, "关键变化", "compact direct row 已超过 MPR diagnostic 的 Acc，并作为主 direct-3D 结果", COLORS["teal"], 15)
    d.card(s, 9.15, 3.08, 3.10, 1.35, "边界说明", "该 row 不读 VPR cache / official RGB SAM；RGB/GrabCut 是轻量 support guard trick", COLORS["amber"], 14)
    d.image(s, FIG_DIR / "lerf_2d3d_ovs_qualitative.png", 0.80, 5.05, 11.7, 1.42, "2D and 3D Open-Vocabulary Query on LERF-OVS")
    d.note("这里要准确讲：compact row 是主线，但 strict no-RGB one-map ablation 另有数值；当前最佳使用轻量 RGB support guard。")

    # 13
    s = d.slide("主结果三：ScanNet VALA-aligned direct point-query", "Results")
    d.table(
        s,
        0.60,
        1.25,
        12.15,
        3.45,
        [
            ["Method", "19 mIoU", "19 mAcc", "15 mIoU", "15 mAcc", "10 mIoU", "10 mAcc"],
            ["LangSplatV2", "14.75", "25.47", "17.09", "35.68", "22.83", "41.52"],
            ["OpenGaussian", "27.73", "42.01", "29.67", "46.15", "39.93", "57.34"],
            ["Dr. Splat", "29.31", "47.68", "33.25", "54.33", "44.19", "65.19"],
            ["VALA", "32.11", "50.05", "35.10", "54.77", "46.21", "65.61"],
            ["GaussFM", "38.06", "61.29", "38.71", "63.15", "47.11", "72.00"],
        ],
        font_size=9,
        emphasize_rows=[5],
    )
    d.image(s, FIG_DIR / "scannet_openvocab_3d_query_qualitative.png", 0.95, 5.05, 11.20, 1.45, "ScanNet open-vocabulary 3D query qualitative")
    d.note("这页只按 VALA-aligned 协议说明对比来源，不出现未发表方法名称。")

    # 14
    s = d.slide("重建 Scene Features vs. 原始 frame-wise RADIO", "Results")
    d.table(
        s,
        0.65,
        1.25,
        7.55,
        3.35,
        [
            ["Task", "Metric", "Frame-wise RADIO", "GaussFM", "Delta"],
            ["LERF text grounding", "mIoU", "0.4634", "0.5707", "+0.1073"],
            ["SAM3 point prompt", "mIoU", "0.3700", "0.4173", "+0.0473"],
            ["SAM3 box prompt", "mIoU", "0.6560", "0.6638", "+0.0079"],
            ["SAM3 mask propagation", "mIoU", "0.3583", "0.3756", "+0.0173"],
            ["DINOv3 dense matching", "score", "0.8547", "0.9048", "+0.0501"],
            ["DINOv3 mask propagation", "mIoU", "0.4606", "0.4677", "+0.0071"],
        ],
        font_size=9,
        emphasize_rows=[1, 2, 3, 4, 5, 6],
    )
    d.card(s, 8.55, 1.45, 3.65, 1.28, "解释口径", "不是“teacher 变弱”，而是多视角 compact memory 对原始逐帧 RADIO 做了重建、去噪和场景一致化。", COLORS["blue"], 14)
    d.image(s, FIG_DIR / "lerf_adaptor_downstream_qualitative.png", 8.40, 3.20, 3.95, 2.25, "Frozen-head tasks qualitative")
    d.note("这里直接支撑用户关心的：student field 全面强于 frame-wise RADIO 的 selected frozen-head tasks。")

    # 15
    s = d.slide("定性一：LERF 2D + 3D Open-Vocabulary Query", "Qualitative")
    d.image(s, FIG_DIR / "lerf_2d3d_ovs_qualitative.png", 0.55, 1.18, 12.25, 5.55)
    d.text(s, 0.85, 6.82, 11.55, 0.25, "2D OVS：RGB 上的 similarity/mask；3D OVS：先选 primitive，再白底渲染 selected Gaussians。", 11, COLORS["muted"])
    d.note("主文定性图应选我们强、baseline 弱但合理的样本；2D baseline 建议 LangSplatV2，3D baseline 建议 Dr. Splat。")

    # 16
    s = d.slide("定性二：ScanNet Open-Vocabulary 3D Query", "Qualitative")
    d.image(s, FIG_DIR / "scannet_openvocab_3d_query_qualitative.png", 0.55, 1.20, 12.25, 5.45)
    d.text(s, 0.86, 6.76, 11.45, 0.25, "主文建议放 binary query point cloud；全类别彩色 point cloud 更适合放 supplementary。", 11, COLORS["muted"])
    d.note("ScanNet 定性不宜太花，重点展示某一开放词汇类别是否被查到。")

    # 17
    s = d.slide("定性三：消融可视化要服务核心贡献", "Qualitative")
    d.image(s, FIG_DIR / "lerf_direct3d_support_policy_ablation_qualitative.png", 0.65, 1.20, 5.85, 3.25, "Direct3D support calibration")
    d.image(s, FIG_DIR / "lerf_rendered_boundary_calibration_qualitative.png", 6.75, 1.20, 5.85, 3.25, "Rendered boundary calibration")
    d.bullets(
        s,
        0.95,
        4.85,
        11.4,
        1.35,
        [
            "最值得保留：Direct3D support calibration，因为它对应小物体漏检、碎片化和 low-visibility support。",
            "Rendered boundary calibration 若可视化差异不明显，可在主文降级，换成更清晰的 feature-memory / support-calibration 消融。",
        ],
        14,
    )
    d.note("这页给导师看图表取舍：不是每个模块都放 qualitative。")

    # 18
    s = d.slide("定量消融：贡献大小排序", "Ablation")
    d.table(
        s,
        0.60,
        1.20,
        12.15,
        4.35,
        [
            ["Contribution", "Task", "Delta", "结论"],
            ["MPR-to-field registration", "Direct3D", "mIoU +0.366 / Acc +0.517", "最大直接收益"],
            ["Foundation-space codec", "Rendered", "LocAcc +0.327 / mIoU +0.225", "核心架构"],
            ["GaussFM vs RADIO", "2D usability", "mIoU +0.126", "主 claim 证据"],
            ["Geometry warm-start", "Rendered", "LocAcc +0.056 / mIoU +0.061", "稳定训练"],
            ["Support calibration", "Direct3D", "mIoU +0.053 / Acc +0.032", "解决小物体/碎片"],
            ["DINO/SAM probes", "2D probes", "selected metrics positive", "辅助支撑"],
        ],
        font_size=9,
        emphasize_rows=[1, 2, 3, 5],
    )
    d.note("这页解决“模块太多”的问题：按贡献排序，不平均讲。")

    # 19
    s = d.slide("Storage / Efficiency：小存储与单次 query latency", "Deployment")
    d.table(
        s,
        0.60,
        1.23,
        6.25,
        3.15,
        [
            ["Scene", "#G", "Direct 1280-D", "Latent", "Saving"],
            ["Figurines", "168k", "412.1 MiB", "20.6 MiB", "20.0x"],
            ["Ramen", "383k", "934.3 MiB", "46.7 MiB", "20.0x"],
            ["Teatime", "460k", "1123.4 MiB", "56.2 MiB", "20.0x"],
            ["Waldo", "692k", "1688.8 MiB", "84.4 MiB", "20.0x"],
        ],
        font_size=9,
        emphasize_rows=[1, 2, 3, 4],
    )
    d.table(
        s,
        7.15,
        1.23,
        5.35,
        2.45,
        [
            ["Task", "ms/query", "Peak VRAM"],
            ["LERF rendered-view", "350.5", "2076 MiB"],
            ["LERF direct 3D", "2367.7", "--"],
            ["ScanNet point query", "387.9", "1666 MiB"],
        ],
        font_size=10,
        emphasize_rows=[1, 2, 3],
    )
    d.bullets(
        s,
        7.20,
        4.10,
        5.10,
        1.45,
        [
            "Feature package 中的 decoder/refiner 是 scene-global，场景越大越容易被摊薄。",
            "后续若做大规模室外场景，storage 优势会比 LERF 小场景更明显。",
        ],
        13,
    )
    d.note("这一页按用户要求突出 single-query latency，而不是总评估耗时。")

    # 20
    s = d.slide("投稿叙事：建议面向 TPAMI 的组织方式", "Writing")
    d.table(
        s,
        0.80,
        1.25,
        11.65,
        4.15,
        [
            ["章节", "主旨", "图表"],
            ["Introduction", "高维 foundation feature 到 compact 3D memory 的问题定义", "Fig. 1 overall"],
            ["Related Work", "3DGS open-vocab / feature distillation / scene memory", "-"],
            ["Method", "geometry scaffold + hybrid compact memory + support calibration", "Fig. 2 details"],
            ["Experiments", "三 benchmark 主表 + feature usability + storage/efficiency", "Tables 1-5"],
            ["Ablations", "按贡献排序，而不是堆模块", "Ablation table + selected qualitative"],
            ["Appendix", "协议、baseline provenance、更多 qualitative/failure cases", "supp figures"],
        ],
        font_size=10,
        emphasize_rows=[1, 3, 4],
    )
    d.note("顶刊不是只看指标，还看叙事是否清楚、协议是否无歧义。")

    # 21
    s = d.slide("当前风险与边界：提前防 reviewer 质疑", "Risk Control")
    d.table(
        s,
        0.75,
        1.25,
        11.75,
        4.40,
        [
            ["风险点", "建议表述"],
            ["DINO/SAM/SigLIP2 是否是真实多教师？", "不是。它们是 RADIO-derived adaptor-space consistency / frozen-head evaluation。"],
            ["RGB/GrabCut 是否破坏 compact claim？", "它是无额外网络、无先验 cache 的轻量 support guard；strict no-RGB row 作为 ablation 保留。"],
            ["VPR cache 是否推理依赖？", "不是。MPR 是训练桥梁，把多视角证据压回 compact field。"],
            ["0.501 direct3D 是否是纯 one-map？", "最佳 row 不读 VPR/SAM，但含 RGB support guard；strict pure one-map row 另报。"],
            ["Storage 为什么 feature package 不够低？", "per-scene global heads 固定开销在小场景明显；随 Gaussian 数增大被摊薄。"],
        ],
        font_size=10,
        emphasize_rows=[1, 2, 3, 4],
    )
    d.note("这页适合导师问答，口径要比论文正文更直白。")

    # 22
    s = d.slide("下一步：把投稿包从“可投”打磨到“强稿”", "Next Steps")
    d.card(s, 0.85, 1.35, 3.55, 1.40, "图表", "重画顶刊级 Fig. 1/2；主文 qualitative 只保留最强样本", COLORS["blue"], 15)
    d.card(s, 4.90, 1.35, 3.55, 1.40, "实验", "冻结 final registry；补 mutual-training 2x2 与 strict compact ablation", COLORS["teal"], 15)
    d.card(s, 8.95, 1.35, 3.55, 1.40, "写作", "按 TPAMI 结构压缩表格；消融按贡献排序；风险边界写清楚", COLORS["amber"], 15)
    d.bullets(
        s,
        1.05,
        3.35,
        11.1,
        2.0,
        [
            "短期目标：导师汇报时先讲清“为什么这不是简单 feature distillation”，再展示三任务结果。",
            "投稿目标：把方法名、贡献点、图表和实验口径统一到 Compact Foundation-Feature Gaussian Memory。",
            "保守状态：论文主线已经基本闭环；最需要加强的是图像质量、术语一致性和少量关键 ablation 的最终同步。",
        ],
        15,
    )
    d.note("收尾时给出明确可执行清单，而不是泛泛说继续优化。")

    d.save()


if __name__ == "__main__":
    build()
