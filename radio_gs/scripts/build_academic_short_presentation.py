#!/usr/bin/env python3
"""Build a concise academic presentation deck for the GaussFM paper."""

from __future__ import annotations

from pathlib import Path
from typing import Iterable, Sequence

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_AUTO_SIZE, MSO_VERTICAL_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


REPO_ROOT = Path(__file__).resolve().parents[2]
PAPER_DIR = REPO_ROOT / "paper"
FIG_DIR = PAPER_DIR / "figures"
OUT_PPTX = PAPER_DIR / "gaussfm_academic_short_presentation.pptx"
OUT_OUTLINE = PAPER_DIR / "gaussfm_academic_short_presentation_outline.md"

SLIDE_W = 13.333
SLIDE_H = 7.5

FONT = "Aptos"
FONT_ALT = "Arial"

COLORS = {
    "ink": RGBColor(24, 31, 42),
    "muted": RGBColor(82, 94, 109),
    "paper": RGBColor(248, 250, 252),
    "panel": RGBColor(255, 255, 255),
    "line": RGBColor(214, 220, 228),
    "teal": RGBColor(18, 126, 126),
    "blue": RGBColor(44, 93, 165),
    "green": RGBColor(68, 135, 93),
    "amber": RGBColor(196, 123, 39),
    "red": RGBColor(174, 68, 68),
    "soft_teal": RGBColor(229, 246, 245),
    "soft_blue": RGBColor(232, 239, 252),
    "soft_amber": RGBColor(252, 244, 230),
}


class ShortDeck:
    def __init__(self) -> None:
        self.prs = Presentation()
        self.prs.slide_width = Inches(SLIDE_W)
        self.prs.slide_height = Inches(SLIDE_H)
        self.blank = self.prs.slide_layouts[6]
        self.outline: list[tuple[str, list[str]]] = []

    def save(self) -> None:
        OUT_PPTX.parent.mkdir(parents=True, exist_ok=True)
        self.prs.save(OUT_PPTX)
        lines = [
            "# GaussFM Academic Short Presentation Outline",
            "",
            f"- PPTX: `{OUT_PPTX}`",
            "- Format: 11-slide concise academic talk deck",
            "- Narrative: problem -> method -> evidence -> takeaway",
            "",
        ]
        for idx, (title, notes) in enumerate(self.outline, 1):
            lines.append(f"## {idx}. {title}")
            for note in notes:
                lines.append(f"- {note}")
            lines.append("")
        OUT_OUTLINE.write_text("\n".join(lines), encoding="utf-8")

    def slide(self, title: str, section: str | None = None) -> object:
        slide = self.prs.slides.add_slide(self.blank)
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = COLORS["paper"]
        if section:
            self.text(slide, 0.55, 0.22, 4.0, 0.22, section.upper(), 8, COLORS["teal"], bold=True)
        self.text(slide, 0.55, 0.47, 11.4, 0.46, title, 23, COLORS["ink"], bold=True)
        self.line(slide, 0.55, 1.06, 12.72, 1.06, COLORS["line"])
        self.footer(slide)
        self.outline.append((title, []))
        return slide

    def note(self, text: str) -> None:
        self.outline[-1][1].append(text)

    def footer(self, slide: object) -> None:
        idx = len(self.prs.slides)
        self.text(slide, 0.55, 7.09, 4.1, 0.18, "GaussFM | Academic Short Talk", 7, COLORS["muted"])
        self.text(slide, 12.20, 7.09, 0.55, 0.18, str(idx), 7, COLORS["muted"], align=PP_ALIGN.RIGHT)

    def text(
        self,
        slide: object,
        x: float,
        y: float,
        w: float,
        h: float,
        value: str,
        size: int,
        color: RGBColor | None = None,
        bold: bool = False,
        align: PP_ALIGN | None = None,
    ) -> object:
        box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = box.text_frame
        tf.clear()
        tf.word_wrap = True
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        tf.margin_left = Inches(0.02)
        tf.margin_right = Inches(0.02)
        tf.margin_top = Inches(0.01)
        tf.margin_bottom = Inches(0.01)
        p = tf.paragraphs[0]
        if align is not None:
            p.alignment = align
        r = p.add_run()
        r.text = value
        r.font.name = FONT
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color or COLORS["ink"]
        return box

    def bullets(
        self,
        slide: object,
        x: float,
        y: float,
        w: float,
        h: float,
        items: Sequence[str],
        size: int = 15,
        color: RGBColor | None = None,
    ) -> None:
        box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = box.text_frame
        tf.clear()
        tf.word_wrap = True
        tf.margin_left = Inches(0.10)
        tf.margin_right = Inches(0.04)
        tf.margin_top = Inches(0.03)
        tf.margin_bottom = Inches(0.02)
        for i, item in enumerate(items):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = item
            p.level = 0
            p.space_after = Pt(7)
            p.font.name = FONT
            p.font.size = Pt(size)
            p.font.color.rgb = color or COLORS["ink"]

    def card(
        self,
        slide: object,
        x: float,
        y: float,
        w: float,
        h: float,
        title: str,
        body: str,
        accent: RGBColor,
        body_size: int = 15,
    ) -> None:
        shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
        shape.fill.solid()
        shape.fill.fore_color.rgb = COLORS["panel"]
        shape.line.color.rgb = COLORS["line"]
        shape.adjustments[0] = 0.04
        bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x), Inches(y), Inches(0.06), Inches(h))
        bar.fill.solid()
        bar.fill.fore_color.rgb = accent
        bar.line.fill.background()
        self.text(slide, x + 0.18, y + 0.14, w - 0.30, 0.26, title, 10, COLORS["muted"], bold=True)
        self.text(slide, x + 0.18, y + 0.48, w - 0.30, h - 0.55, body, body_size, COLORS["ink"], bold=True)

    def chip(self, slide: object, x: float, y: float, w: float, label: str, color: RGBColor) -> None:
        shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(0.34))
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()
        shape.adjustments[0] = 0.5
        tf = shape.text_frame
        tf.clear()
        tf.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = label
        r.font.name = FONT
        r.font.size = Pt(8.5)
        r.font.bold = True
        r.font.color.rgb = RGBColor(255, 255, 255)

    def line(self, slide: object, x1: float, y1: float, x2: float, y2: float, color: RGBColor) -> None:
        line = slide.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
        line.line.color.rgb = color
        line.line.width = Pt(1.0)

    def image(self, slide: object, path: Path, x: float, y: float, w: float, h: float) -> None:
        if not path.exists():
            self.card(slide, x, y, w, h, "Missing figure", str(path.relative_to(REPO_ROOT)), COLORS["red"], 11)
            return
        with Image.open(path) as im:
            iw, ih = im.size
        box_ratio = w / h
        img_ratio = iw / ih
        if img_ratio >= box_ratio:
            draw_w = w
            draw_h = w / img_ratio
            draw_x = x
            draw_y = y + (h - draw_h) / 2
        else:
            draw_h = h
            draw_w = h * img_ratio
            draw_x = x + (w - draw_w) / 2
            draw_y = y
        slide.shapes.add_picture(str(path), Inches(draw_x), Inches(draw_y), Inches(draw_w), Inches(draw_h))

    def table(
        self,
        slide: object,
        x: float,
        y: float,
        w: float,
        h: float,
        data: Sequence[Sequence[str]],
        font_size: int = 10,
        emphasize_rows: Iterable[int] = (),
    ) -> None:
        rows, cols = len(data), len(data[0])
        shape = slide.shapes.add_table(rows, cols, Inches(x), Inches(y), Inches(w), Inches(h))
        table = shape.table
        emph = set(emphasize_rows)
        for i, row in enumerate(data):
            for j, value in enumerate(row):
                cell = table.cell(i, j)
                cell.text = value
                cell.margin_left = Inches(0.04)
                cell.margin_right = Inches(0.04)
                cell.margin_top = Inches(0.02)
                cell.margin_bottom = Inches(0.02)
                fill = COLORS["soft_blue"] if i == 0 else (COLORS["soft_teal"] if i in emph else COLORS["panel"])
                cell.fill.solid()
                cell.fill.fore_color.rgb = fill
                for p in cell.text_frame.paragraphs:
                    p.alignment = PP_ALIGN.CENTER if j else PP_ALIGN.LEFT
                    for r in p.runs:
                        r.font.name = FONT
                        r.font.size = Pt(font_size)
                        r.font.bold = i == 0 or i in emph
                        r.font.color.rgb = COLORS["ink"]


def build_deck() -> ShortDeck:
    d = ShortDeck()

    # 1
    s = d.slide("GaussFM: Compact Foundation-Feature Gaussian Memory", "thesis")
    d.text(
        s,
        0.75,
        1.42,
        11.8,
        0.75,
        "Turn view-local foundation features into a compact, queryable 3D scene memory.",
        25,
        COLORS["ink"],
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    d.card(s, 1.05, 2.55, 3.35, 1.25, "Input", "posed RGB + frozen foundation features", COLORS["blue"], 15)
    d.card(s, 4.98, 2.55, 3.35, 1.25, "Memory", "low-dimensional Gaussian feature codes", COLORS["teal"], 15)
    d.card(s, 8.90, 2.55, 3.35, 1.25, "Queries", "2D rendered maps, 3D primitives, point queries", COLORS["green"], 15)
    d.text(s, 1.30, 4.45, 10.7, 0.52, "One representation; three open-vocabulary query modes.", 20, COLORS["teal"], bold=True, align=PP_ALIGN.CENTER)
    d.note("Open with one sentence: GaussFM is a deployable 3D foundation-feature memory, not a task-specific classifier.")

    # 2
    s = d.slide("Why 2D foundation features are not yet 3D scene memory", "motivation")
    d.card(s, 0.80, 1.45, 3.75, 1.40, "View-local", "single-frame features vary with viewpoint, scale, and occlusion", COLORS["blue"], 14)
    d.card(s, 4.80, 1.45, 3.75, 1.40, "High-dimensional", "storing dense 1280-D features per view or primitive is expensive", COLORS["amber"], 14)
    d.card(s, 8.80, 1.45, 3.75, 1.40, "Not directly 3D", "2D heatmaps do not automatically answer primitive- or point-level queries", COLORS["red"], 14)
    d.text(s, 0.90, 3.58, 11.6, 0.45, "The missing object is a compact scene-level foundation-feature memory.", 20, COLORS["ink"], bold=True, align=PP_ALIGN.CENTER)
    d.bullets(
        s,
        1.55,
        4.35,
        10.4,
        1.25,
        [
            "It should render dense feature maps for novel views.",
            "It should support direct Gaussian/point queries without a view cache.",
            "It should preserve downstream usability under frozen heads.",
        ],
        15,
        COLORS["muted"],
    )
    d.note("Motivation slide: identify the representation gap rather than criticizing prior work broadly.")

    # 3
    s = d.slide("Research question", "problem")
    d.text(
        s,
        1.05,
        1.55,
        11.3,
        1.0,
        "How can high-dimensional, multiview foundation features be compressed into a 3D Gaussian field without losing queryability?",
        25,
        COLORS["ink"],
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    d.card(s, 1.00, 3.20, 3.35, 1.38, "Q1. Compression", "low memory, preserved RADIO-compatible feature space", COLORS["teal"], 14)
    d.card(s, 4.98, 3.20, 3.35, 1.38, "Q2. Consistency", "multiview evidence must become stable primitive support", COLORS["blue"], 14)
    d.card(s, 8.95, 3.20, 3.35, 1.38, "Q3. Query modes", "one field should serve rendered, primitive, and point queries", COLORS["green"], 14)
    d.note("This slide creates the slots that the method components will fill.")

    # 4
    s = d.slide("Method overview: one compact memory, three query modes", "method")
    d.image(s, FIG_DIR / "figure1_overall_framework.png", 0.85, 1.28, 11.85, 4.85)
    d.text(
        s,
        0.90,
        6.20,
        11.6,
        0.38,
        "Training reconstructs RADIO-compatible scene features; inference queries the stored Gaussian memory.",
        15,
        COLORS["muted"],
        align=PP_ALIGN.CENTER,
    )
    d.note("Use the framework figure as the visual anchor. Speak in train/inference terms, not module-by-module detail.")

    # 5
    s = d.slide("Core components in the compact feature field", "method")
    d.image(s, FIG_DIR / "figure2_method_details.png", 0.65, 1.20, 6.15, 4.75)
    d.card(s, 7.10, 1.34, 5.25, 0.82, "1. Compact Gaussian codes", "store low-dimensional feature latents per Gaussian", COLORS["teal"], 13)
    d.card(s, 7.10, 2.30, 5.25, 0.82, "2. Spatial context", "inject neighborhood and scene-level support for stable features", COLORS["blue"], 13)
    d.card(s, 7.10, 3.26, 5.25, 0.82, "3. Foundation-space decoding", "recover RADIO-compatible features on demand", COLORS["green"], 13)
    d.card(s, 7.10, 4.22, 5.25, 0.82, "4. Reliability / visibility", "calibrate support for rendered and direct queries", COLORS["amber"], 13)
    d.text(s, 7.15, 5.52, 5.10, 0.50, "Compact storage is paired with reconstruction and support calibration.", 15, COLORS["ink"], bold=True)
    d.note("Method slide stays concise: four components, each tied to a bottleneck.")

    # 6
    s = d.slide("Main contributions and innovations", "contributions")
    d.card(s, 0.82, 1.38, 3.70, 2.18, "Contribution 1", "A compact foundation-feature Gaussian memory that stores low-dimensional codes rather than full RADIO features.", COLORS["teal"], 15)
    d.card(s, 4.82, 1.38, 3.70, 2.18, "Contribution 2", "A unified query interface for rendered 2D grounding, direct 3D object selection, and point-level open-vocabulary query.", COLORS["blue"], 15)
    d.card(s, 8.82, 1.38, 3.70, 2.18, "Contribution 3", "Evidence that reconstructed scene features can remain compact while improving selected downstream usability.", COLORS["green"], 15)
    d.text(s, 1.05, 4.30, 11.2, 0.42, "The novelty is the reusable scene memory, not a benchmark-specific readout.", 19, COLORS["ink"], bold=True, align=PP_ALIGN.CENTER)
    d.note("This is the explicit innovation slide requested by the user.")

    # 7
    s = d.slide("Result 1: rendered-view open-vocabulary grounding", "evidence")
    d.table(
        s,
        0.88,
        1.38,
        5.45,
        2.55,
        [
            ["Method", "mIoU", "Acc"],
            ["LangSplat", "33.92", "64.01"],
            ["LangSplatV2", "46.24", "75.84"],
            ["OpenGaussian", "49.74", "72.75"],
            ["GaussFM", "64.98", "82.68"],
        ],
        11,
        emphasize_rows=[4],
    )
    d.image(s, FIG_DIR / "lerf_2d3d_ovs_qualitative.png", 6.75, 1.30, 5.75, 4.70)
    d.text(s, 0.95, 4.45, 5.30, 0.75, "Rendered GaussFM features support strong 2D open-vocabulary localization under the same reproduced protocol.", 17, COLORS["ink"], bold=True)
    d.note("Evidence 1: show that the memory renders useful dense maps.")

    # 8
    s = d.slide("Result 2: direct 3D open-vocabulary object selection", "evidence")
    d.table(
        s,
        0.82,
        1.38,
        5.75,
        2.40,
        [
            ["Method", "mIoU", "Acc"],
            ["OpenGaussian", "41.06", "51.44"],
            ["Dr. Splat", "39.77", "65.48"],
            ["LangSplatV2", "35.87", "55.80"],
            ["GaussFM", "54.36", "80.84"],
        ],
        11,
        emphasize_rows=[4],
    )
    d.card(s, 0.95, 4.25, 5.35, 1.08, "Interpretation", "The compact memory can be queried at the primitive level, not only after rendering a 2D heatmap.", COLORS["blue"], 15)
    d.image(s, FIG_DIR / "lerf_direct3d_support_policy_ablation_qualitative.png", 6.88, 1.25, 5.55, 4.78)
    d.note("Evidence 2: direct 3D query validates the stored scene memory.")

    # 9
    s = d.slide("Result 3: VALA-aligned ScanNet point query", "evidence")
    d.table(
        s,
        0.72,
        1.38,
        7.05,
        3.08,
        [
            ["Method", "19 mIoU", "15 mIoU", "10 mIoU"],
            ["OpenGaussian", "27.73", "29.67", "39.93"],
            ["Dr. Splat", "29.31", "33.25", "44.19"],
            ["OccamLGS", "31.93", "34.25", "45.16"],
            ["VALA", "32.11", "35.10", "46.21"],
            ["GaussFM", "36.55", "42.78", "57.85"],
        ],
        10,
        emphasize_rows=[5],
    )
    d.image(s, FIG_DIR / "scannet_openvocab_3d_query_qualitative.png", 8.05, 1.35, 4.75, 3.72)
    d.text(s, 0.90, 5.05, 11.8, 0.45, "The same feature-memory idea transfers beyond LERF scenes to point-level open-vocabulary querying.", 17, COLORS["ink"], bold=True, align=PP_ALIGN.CENTER)
    d.note("Evidence 3: ScanNet establishes cross-dataset point-query transfer under a VALA-aligned protocol.")

    # 10
    s = d.slide("Why it works: compact reconstruction plus support calibration", "analysis")
    d.table(
        s,
        0.78,
        1.35,
        7.10,
        3.25,
        [
            ["Ablation", "Task", "Main gain"],
            ["Foundation-space reconstruction", "LERF rendered", "+0.366 mIoU"],
            ["Frame-wise RADIO -> GaussFM", "same evaluator", "+0.107 mIoU"],
            ["Prompt ensemble + component support", "Direct 3D", "+0.053 mIoU"],
            ["SAM3 point-prompt feature readout", "frozen-head probe", "+0.047 mIoU"],
        ],
        10,
        emphasize_rows=[1, 2],
    )
    d.card(s, 8.25, 1.45, 4.35, 1.18, "Compactness", "20x smaller latent payload than explicit 1280-D fp16 features.", COLORS["teal"], 15)
    d.card(s, 8.25, 2.95, 4.35, 1.18, "Usability", "6/6 selected primary downstream probes favor rendered GaussFM features.", COLORS["green"], 15)
    d.card(s, 8.25, 4.45, 4.35, 1.18, "Boundary", "Secondary metrics and diagnostics are reported separately to avoid overclaiming.", COLORS["amber"], 15)
    d.note("This slide makes the method credible without flooding the audience with all ablation rows.")

    # 11
    s = d.slide("Takeaway", "summary")
    d.text(
        s,
        0.95,
        1.42,
        11.5,
        0.92,
        "GaussFM converts frame-wise foundation features into a compact, reusable, queryable 3D Gaussian scene memory.",
        26,
        COLORS["ink"],
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    d.card(s, 1.00, 3.00, 3.45, 1.22, "Compact", "low-dimensional per-Gaussian feature memory", COLORS["teal"], 15)
    d.card(s, 4.95, 3.00, 3.45, 1.22, "Unified", "2D rendered, 3D primitive, and point queries", COLORS["blue"], 15)
    d.card(s, 8.90, 3.00, 3.45, 1.22, "Effective", "strong reproduced results across three protocols", COLORS["green"], 15)
    d.text(s, 1.15, 5.25, 11.0, 0.48, "One compact memory can carry foundation-model semantics into 3D scene understanding.", 19, COLORS["teal"], bold=True, align=PP_ALIGN.CENTER)
    d.note("Close by returning to the central thesis and the three evidence tracks.")

    return d


def main() -> None:
    deck = build_deck()
    deck.save()
    print(f"Wrote {OUT_PPTX}")
    print(f"Wrote {OUT_OUTLINE}")


if __name__ == "__main__":
    main()
